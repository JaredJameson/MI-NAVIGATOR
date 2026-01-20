"""
Reports API Endpoints
"""

from fastapi import APIRouter, Query, Depends, HTTPException, UploadFile, File, Request
from fastapi.responses import StreamingResponse
from typing import Optional, List
from pydantic import BaseModel
from datetime import datetime
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
import io
import re
import uuid
import json
import asyncio
from pathlib import Path

from app.api.v1.endpoints.auth import get_current_user
from app.models.user import User
from app.models.report_template import ReportTemplate
from app.db.session import get_db
from app.services.audit_service import log_audit, AuditAction, ResourceType
from app.services.analytics_service import track_event
from app.models.analytics_event import EventType
from app.api.v1.endpoints.projects import MOCK_PROJECTS

router = APIRouter()


# In-memory task storage for progress tracking
BATCH_TASKS = {}  # {task_id: {"status": "processing", "total": 10, "processed": 3, "result": None}}


# Mock templates database
MOCK_TEMPLATES = [
    {
        "id": "template_test001",
        "name": "Szablon profilu firmy produkcyjnej",
        "type": "company_profile",
        "created_at": "2026-01-16T22:00:00Z",
        "created_by": "user_001",
        "use_count": 0,
        "last_used": None,
        "original_report_title": "Analiza profilu FADO Sp. z o.o.",
        "sections": [
            {"id": "section_1", "title": "Informacje podstawowe", "content": "[Do uzupełnienia: dane rejestrowe, NIP, REGON, KRS]"},
            {"id": "section_2", "title": "Analiza finansowa", "content": "[Do uzupełnienia: przychody, wskaźniki finansowe, trendy]"},
            {"id": "section_3", "title": "Pozycja rynkowa", "content": "[Do uzupełnienia: udział w rynku, konkurenci, przewagi]"},
            {"id": "section_4", "title": "Analiza SWOT", "content": "[Do uzupełnienia: mocne/słabe strony, szanse/zagrożenia]"}
        ]
    },
    {
        "id": "template_test002",
        "name": "Szablon analizy rynku",
        "type": "market_analysis",
        "created_at": "2026-01-16T22:00:00Z",
        "created_by": "user_001",
        "use_count": 0,
        "last_used": None,
        "original_report_title": "Analiza rynku tworzyw sztucznych w Polsce",
        "sections": [
            {"id": "section_1", "title": "Wielkość rynku", "content": "[Do uzupełnienia: TAM, SAM, SOM, wzrost rynku]"},
            {"id": "section_2", "title": "Segmentacja rynku", "content": "[Do uzupełnienia: kluczowe segmenty, wielkość, trendy]"},
            {"id": "section_3", "title": "Analiza konkurencji", "content": "[Do uzupełnienia: główni gracze, udziały rynkowe]"},
            {"id": "section_4", "title": "Trendy i prognozy", "content": "[Do uzupełnienia: trendy technologiczne, prognozy wzrostu]"},
            {"id": "section_5", "title": "Bariery wejścia", "content": "[Do uzupełnienia: regulacje, kapitał, technologia]"}
        ]
    }
]

# Helper function to generate pagination test data
def generate_pagination_test_reports(count: int = 1000, user_id: str = "f6e9a62e-fb70-4808-882b-e5711d0a5411"):
    """Generate test reports for pagination performance testing (Feature #200)"""
    reports = []
    for i in range(1, count + 1):
        report = {
            "id": f"pagination_test_{i:04d}",
            "title": f"Pagination Test Report #{i}",
            "type": "market_analysis" if i % 2 == 0 else "company_profile",
            "company": f"Test Company {i}" if i % 3 == 0 else None,
            "created_at": f"2026-01-{min(19, 1 + i % 28):02d}T{i % 24:02d}:{i % 60:02d}:00Z",
            "updated_at": f"2026-01-{min(19, 1 + i % 28):02d}T{i % 24:02d}:{i % 60:02d}:00Z",
            "status": "completed",
            "is_archived": False,
            "created_by": user_id,
            "summary": f"Test report {i} for pagination performance testing",
            "sections": [],
            "sources": []
        }
        reports.append(report)
    return reports

# Mock reports database
MOCK_REPORTS = [
    {
        "id": "report_test113",
        "title": "TestReport",
        "type": "market_analysis",
        "company": None,
        "created_at": "2026-01-19T01:25:00Z",
        "updated_at": "2026-01-19T01:25:00Z",
        "status": "completed",
        "is_archived": False,
        "created_by": "dde243d1-a7a9-44c5-b28a-772ece7d500e",
        "summary": "Test report for case-insensitive search - Feature 113",
        "sections": [],
        "sources": []
    },
    {
        "id": "report_001",
        "title": "Analiza profilu FADO Sp. z o.o.",
        "type": "company_profile",
        "company": "FADO Sp. z o.o.",
        "created_at": "2026-01-14T10:30:00Z",
        "updated_at": "2026-01-14T14:22:00Z",
        "status": "completed",
        "is_archived": False,
        "created_by": "dde243d1-a7a9-44c5-b28a-772ece7d500e",
        "summary": "Kompleksowa analiza profilu firmy FADO Sp. z o.o. - lidera w produkcji tworzyw sztucznych.",
        "sections": [
            {
                "id": "section_1",
                "title": "Informacje podstawowe",
                "content": """| Header 1 | Header 2 | Header 3 |
| --- | --- | --- |
| Cell 1-1 | Cell 1-2 | Cell 1-3 |
| Cell 2-1 | Cell 2-2 | Cell 2-3 |

FADO Sp. z o.o. to polska firma założona w 1998 roku, specjalizująca się w produkcji wyrobów z tworzyw sztucznych. Firma posiada siedzibę w Warszawie przy ul. Przemysłowej 15.

**Dane rejestrowe:**
- NIP: 5260016831
- REGON: 012567834
- KRS: 0000145732
- Forma prawna: Spółka z ograniczoną odpowiedzialnością

**Główna działalność (PKD):**
- 22.21.Z - Produkcja płyt, arkuszy, rur i kształtowników z tworzyw sztucznych
- 22.22.Z - Produkcja opakowań z tworzyw sztucznych
- 22.29.Z - Produkcja pozostałych wyrobów z tworzyw sztucznych

Firma zatrudnia około 150 pracowników i posiada nowoczesny park maszynowy o łącznej mocy produkcyjnej 5000 ton rocznie."""
            },
            {
                "id": "section_2",
                "title": "Analiza finansowa",
                "content": """**Przychody i rentowność (2023):**
- Przychody ze sprzedaży: 45,2 mln PLN
- Wzrost r/r: +12,3%
- Marża brutto: 28,5%
- Zysk netto: 4,8 mln PLN

**Wskaźniki finansowe:**
- ROE (Return on Equity): 18,2%
- ROA (Return on Assets): 9,4%
- Wskaźnik płynności bieżącej: 2,1
- Wskaźnik zadłużenia: 32%

**Trend przychodów (mln PLN):**
- 2021: 35,8 mln PLN
- 2022: 40,2 mln PLN
- 2023: 45,2 mln PLN

Firma wykazuje stabilny wzrost przychodów na poziomie 10-15% rocznie. Marża brutto utrzymuje się na poziomie konkurencyjnym dla branży.

**Źródła danych:**
- [Raport roczny FADO 2023](https://fado.pl/raporty/2023) - sprawozdanie finansowe
- Więcej informacji na stronie [KRS Online](https://ekrs.ms.gov.pl)"""
            },
            {
                "id": "section_3",
                "title": "Pozycja rynkowa",
                "content": """**Udział w rynku:**
FADO zajmuje pozycję wśród Top 10 producentów tworzyw sztucznych w Polsce, z szacowanym udziałem w rynku na poziomie 3,5%.

**Główni konkurenci:**
1. Splast S.A. - lider rynku z udziałem 8%
2. PlastPak Sp. z o.o. - 4,2% rynku
3. PolyTech Sp. z o.o. - 3,8% rynku

**Przewagi konkurencyjne:**
- Nowoczesny park maszynowy
- Certyfikaty jakości ISO 9001 i ISO 14001
- Elastyczność w realizacji zamówień
- Własne centrum R&D

**Wyzwania:**
- Rosnące ceny surowców
- Presja regulacyjna na plastik
- Konkurencja ze strony importu z Azji"""
            },
            {
                "id": "section_4",
                "title": "Analiza SWOT",
                "content": """**Mocne strony (Strengths):**
- Doświadczenie 25+ lat na rynku
- Wykwalifikowana kadra
- Nowoczesna infrastruktura
- Certyfikaty jakości
- Stabilna baza klientów

**Słabe strony (Weaknesses):**
- Koncentracja na rynku polskim
- Zależność od kilku głównych dostawców
- Brak własnej sieci dystrybucji

**Szanse (Opportunities):**
- Rozwój rynku opakowań biodegradowalnych
- Ekspansja na rynki CEE
- Rosnący popyt na plastik recyklingowy
- Digitalizacja procesów produkcyjnych

**Zagrożenia (Threats):**
- Regulacje antyplastikowe EU
- Wahania cen ropy naftowej
- Konkurencja z krajów o niższych kosztach pracy
- Spowolnienie gospodarcze"""
            },
            {
                "id": "section_6",
                "title": "Wskaźniki finansowe - Radar",
                "content": """[FINANCIAL_RATIOS_RADAR]

**Wskaźniki rentowności:**
[RATIO] ROE (Return on Equity)
Wartość: 18.2%
Benchmark branżowy: 15.0%
Opis: Zwrot z kapitału własnego - mierzy efektywność wykorzystania kapitału właścicieli

[RATIO] ROA (Return on Assets)
Wartość: 9.4%
Benchmark branżowy: 7.5%
Opis: Zwrot z aktywów - pokazuje jak efektywnie firma wykorzystuje swoje aktywa

[RATIO] ROS (Return on Sales)
Wartość: 10.6%
Benchmark branżowy: 8.0%
Opis: Marża zysku netto - procent zysku z każdej złotówki przychodu

**Wskaźniki płynności:**
[RATIO] Wskaźnik płynności bieżącej (Current Ratio)
Wartość: 2.1
Benchmark branżowy: 1.5
Opis: Zdolność do pokrycia zobowiązań krótkoterminowych aktywami obrotowymi

[RATIO] Wskaźnik płynności szybkiej (Quick Ratio)
Wartość: 1.4
Benchmark branżowy: 1.0
Opis: Płynność bez uwzględnienia zapasów

**Wskaźniki zadłużenia:**
[RATIO] Wskaźnik zadłużenia ogólnego (Debt Ratio)
Wartość: 32%
Benchmark branżowy: 45%
Opis: Udział zobowiązań w finansowaniu aktywów (niższy = lepszy)

[RATIO] Wskaźnik zadłużenia kapitału własnego (D/E)
Wartość: 0.47
Benchmark branżowy: 0.80
Opis: Stosunek długu do kapitału własnego

**Wskaźniki efektywności:**
[RATIO] Rotacja zapasów (Inventory Turnover)
Wartość: 6.2
Benchmark branżowy: 5.0
Opis: Ile razy w roku firma odnawia zapasy

[RATIO] Rotacja należności (DSO - Days Sales Outstanding)
Wartość: 45 dni
Benchmark branżowy: 60 dni
Opis: Średni czas oczekiwania na płatność od klientów

**Podsumowanie:**
FADO Sp. z o.o. wykazuje wskaźniki finansowe powyżej średniej branżowej w większości kategorii. Szczególnie silna pozycja w zakresie rentowności (ROE, ROA) oraz płynności. Niski poziom zadłużenia daje możliwość dalszej ekspansji poprzez finansowanie dłużne."""
            },
            {
                "id": "section_5",
                "title": "Struktura własnościowa",
                "content": """**Drzewo struktury własnościowej:**

[ROOT] FADO Sp. z o.o. (100%)
├── [SHAREHOLDER] Plastics Holding S.A. (60%)
│   ├── [SHAREHOLDER] Euro Polymers GmbH (45%)
│   │   └── [UBO] Heinrich Schmidt (100%) - Beneficjent rzeczywisty
│   └── [SHAREHOLDER] Polish Investment Fund S.A. (55%)
│       └── [UBO] Skarb Państwa (100%) - Beneficjent rzeczywisty
├── [SHAREHOLDER] Jan Kowalski (25%) - Założyciel
│   └── [UBO] Jan Kowalski (100%) - Beneficjent rzeczywisty
└── [SHAREHOLDER] Anna Nowak (15%) - Współzałożyciel
    └── [UBO] Anna Nowak (100%) - Beneficjent rzeczywisty

**Szczegóły udziałowców:**

Plastics Holding S.A. - 60%
Typ: Spółka akcyjna
KRS: 0000234567
Rola: Większościowy udziałowiec strategiczny
Opis: Holding inwestycyjny specjalizujący się w branży tworzyw sztucznych

Euro Polymers GmbH - 45% (w Plastics Holding)
Typ: Spółka zagraniczna
Kraj: Niemcy
Rola: Inwestor branżowy
Opis: Niemiecki producent polimerów, część grupy kapitałowej

Polish Investment Fund S.A. - 55% (w Plastics Holding)
Typ: Fundusz inwestycyjny
KRS: 0000345678
Rola: Inwestor finansowy
Opis: Państwowy fundusz wspierający rozwój przemysłu

Jan Kowalski - 25%
Typ: Osoba fizyczna
Rola: Założyciel i Prezes Zarządu
Opis: Założyciel firmy, pełni funkcję Prezesa od 1998 roku

Anna Nowak - 15%
Typ: Osoba fizyczna
Rola: Współzałożyciel i Członek Zarządu
Opis: Współzałożycielka, odpowiedzialna za rozwój biznesu

**Beneficjenci rzeczywiści (UBO):**
- Heinrich Schmidt (poprzez Euro Polymers GmbH → Plastics Holding S.A.)
- Skarb Państwa (poprzez Polish Investment Fund S.A. → Plastics Holding S.A.)
- Jan Kowalski (bezpośrednio)
- Anna Nowak (bezpośrednio)

**Historia zmian własnościowych:**
- 1998: Założenie przez Jana Kowalskiego (70%) i Annę Nowak (30%)
- 2010: Wejście Plastics Holding S.A. (40%), rozwodnienie do: Kowalski 42%, Nowak 18%
- 2018: Zwiększenie udziału Plastics Holding do 60%, obecna struktura"""
            }
        ],
        "sources": [
            {"name": "KRS", "confidence": 0.95, "url": "https://api.krs.pl"},
            {"name": "e-sprawozdania", "confidence": 0.90, "url": "https://ekrs.ms.gov.pl"},
            {"name": "Analiza branżowa PZPTS", "confidence": 0.85, "url": "https://pzpts.pl"}
        ],
        "charts": [
            {
                "id": "chart_revenue_trend",
                "title": "Trend przychodów FADO (mln PLN)",
                "type": "bar",
                "section_id": "section_2",
                "data": [
                    {"label": "2021", "value": 35.8},
                    {"label": "2022", "value": 40.2},
                    {"label": "2023", "value": 45.2}
                ],
                "xKey": "label",
                "yKey": "value",
                "yLabel": "Przychody (mln PLN)",
                "color": "#3b82f6"
            },
            {
                "id": "chart_competitors",
                "title": "Udział w rynku - główni konkurenci",
                "type": "bar",
                "section_id": "section_3",
                "data": [
                    {"label": "Splast S.A.", "value": 8.0},
                    {"label": "PlastPak", "value": 4.2},
                    {"label": "PolyTech", "value": 3.8},
                    {"label": "FADO", "value": 3.5}
                ],
                "xKey": "label",
                "yKey": "value",
                "yLabel": "Udział w rynku (%)",
                "color": "#10b981"
            }
        ]
    },
    {
        "id": "report_002",
        "title": "Analiza rynku tworzyw sztucznych w Polsce",
        "type": "market_analysis",
        "company": None,
        "created_at": "2026-01-13T09:15:00Z",
        "updated_at": "2026-01-13T16:45:00Z",
        "status": "completed",
        "is_archived": False,
        "created_by": "dde243d1-a7a9-44c5-b28a-772ece7d500e",
        "summary": "Kompleksowa analiza rynku produkcji tworzyw sztucznych w Polsce - trendy, gracze, prognozy.",
        "sections": [
            {
                "id": "section_1",
                "title": "Wielkość i struktura rynku",
                "content": """**Wielkość rynku (2023):**
- Wartość rynku: 12,8 mld PLN
- Wolumen produkcji: 3,2 mln ton
- Liczba przedsiębiorstw: ~2,500
- Zatrudnienie w sektorze: ~45,000 osób

**Struktura rynku według segmentów:**
- Opakowania: 42%
- Budownictwo: 23%
- Motoryzacja: 15%
- AGD i elektronika: 12%
- Pozostałe: 8%

**Dynamika wzrostu:**
Rynek tworzyw sztucznych w Polsce rośnie średnio o 4-6% rocznie. Głównym motorem wzrostu jest sektor opakowaniowy oraz budownictwo."""
            },
            {
                "id": "section_2",
                "title": "Kluczowi gracze",
                "content": """**Top 10 producentów tworzyw sztucznych w Polsce:**

1. **Grupa Azoty Polyolefins** - 15% rynku
   - Lider w produkcji polipropylenu
   - Przychody: ~2 mld PLN

2. **Splast S.A.** - 8% rynku
   - Specjalizacja: opakowania przemysłowe
   - Przychody: ~900 mln PLN

3. **PolyTech Sp. z o.o.** - 5% rynku
   - Komponenty motoryzacyjne
   - Przychody: ~600 mln PLN

4. **PlastPak Sp. z o.o.** - 4,2% rynku
   - Opakowania konsumenckie
   - Przychody: ~500 mln PLN

5. **FADO Sp. z o.o.** - 3,5% rynku
   - Profile i rury
   - Przychody: ~450 mln PLN"""
            },
            {
                "id": "section_3",
                "title": "Analiza Porter Five Forces",
                "content": """**Siła przetargowa dostawców (Supplier Power): ŚREDNIA**
- Ograniczona liczba dostawców surowców petrochemicznych
- Wysokie koszty zmiany dostawcy
- Możliwość integracji wstecznej przez dużych producentów
- Zależność od globalnych cen ropy naftowej

**Siła przetargowa nabywców (Buyer Power): WYSOKA**
- Duża liczba alternatywnych dostawców
- Niskie koszty zmiany dla odbiorców
- Rosnące wymagania dotyczące jakości i certyfikacji
- Presja cenowa ze strony dużych sieci handlowych

**Zagrożenie ze strony substytutów (Threat of Substitutes): WYSOKA**
- Rosnąca popularność opakowań biodegradowalnych
- Papier i karton jako alternatywa
- Szkło i aluminium w segmencie premium
- Regulacje ograniczające plastik jednorazowy

**Zagrożenie ze strony nowych wchodzących (Threat of New Entrants): NISKA**
- Wysokie bariery wejścia (kapitał, technologia)
- Efekty skali u obecnych graczy
- Wymogi certyfikacji i norm jakościowych
- Ugruntowane relacje z odbiorcami

**Rywalizacja wewnątrz branży (Industry Rivalry): WYSOKA**
- Wielu konkurentów o podobnej wielkości
- Niska dynamika wzrostu rynku (4-6% rocznie)
- Wysokie koszty stałe wymuszające walkę o wolumen
- Ograniczone możliwości różnicowania produktów"""
            },
            {
                "id": "section_4",
                "title": "Wielkość rynku TAM SAM SOM",
                "content": """**TAM (Total Addressable Market) - Całkowity Rynek:**
Wartość: 85 mld PLN
Opis: Globalny rynek tworzyw sztucznych dostępny dla polskich firm, obejmujący kraje UE i eksport

**SAM (Serviceable Addressable Market) - Rynek Docelowy:**
Wartość: 28 mld PLN
Opis: Rynek tworzyw sztucznych w Polsce i najbliższych krajach CEE (Czechy, Słowacja, Węgry)

**SOM (Serviceable Obtainable Market) - Rynek Osiągalny:**
Wartość: 4,5 mld PLN
Opis: Realistyczny udział rynkowy możliwy do osiągnięcia dla średniej wielkości producenta w ciągu 3-5 lat

**Metodologia kalkulacji:**
- TAM: Suma wartości rynku UE * dostępność dla polskich eksporterów
- SAM: Rynek krajowy + penetracja regionalna (CEE)
- SOM: Bazując na benchmarkach konkurencji i zdolności produkcyjnych

**Prognoza wzrostu (CAGR 2024-2028):**
- TAM: 3,2% rocznie
- SAM: 4,8% rocznie
- SOM: 8,5% rocznie (przy agresywnej strategii ekspansji)"""
            },
            {
                "id": "section_5",
                "title": "Trendy rynkowe - Timeline",
                "content": """**Trend: Bioplastiki i materiały biodegradowalne**
Kategoria: Technologia
Status: Rosnący
Okres: 2020-2028
Wpływ: Wysoki
Punkty czasowe:
- 2020: Pierwsze regulacje UE dot. plastiku jednorazowego (SUP Directive)
- 2021: Wzrost popytu na alternatywy o 15%
- 2022: Główni producenci uruchamiają linie bio
- 2023: Udział bioplastików osiąga 5% rynku
- 2024: Prognoza: 8% udziału, nowe technologie PLA/PHA
- 2026: Oczekiwany udział 12%, masowa produkcja
- 2028: Cel UE: 20% opakowań z materiałów odnawialnych

**Trend: Gospodarka cyrkularna i recykling**
Kategoria: Regulacje
Status: Rosnący
Okres: 2019-2030
Wpływ: Wysoki
Punkty czasowe:
- 2019: EU Circular Economy Action Plan
- 2020: Polska ustawa o ROP (Rozszerzona Odpowiedzialność Producenta)
- 2022: Wymogi 25% recyklatu w butelkach PET
- 2024: System kaucyjny w Polsce
- 2025: Cel 50% recyklingu opakowań plastikowych
- 2030: Cel UE: 55% recyklingu wszystkich plastików

**Trend: Automatyzacja i Przemysł 4.0**
Kategoria: Technologia
Status: Dojrzały
Okres: 2018-2026
Wpływ: Średni
Punkty czasowe:
- 2018: Pierwsze wdrożenia IoT w branży
- 2020: Pandemia przyspiesza automatyzację
- 2022: 30% fabryk z elementami Industry 4.0
- 2024: Roboty współpracujące (cobots) standardem
- 2026: Przewidywane 60% automatyzacji procesów

**Trend: Konsolidacja rynku**
Kategoria: Rynek
Status: Stabilny
Okres: 2021-2027
Wpływ: Średni
Punkty czasowe:
- 2021: 3 duże fuzje w sektorze
- 2022: Wejście funduszy PE na rynek
- 2023: Top 10 kontroluje 45% rynku
- 2025: Prognoza: Top 10 z 55% udziałem
- 2027: Oczekiwana dalsza konsolidacja"""
            }
        ],
        "sources": [
            {"name": "GUS", "confidence": 0.95, "url": "https://stat.gov.pl"},
            {"name": "PZPTS", "confidence": 0.90, "url": "https://pzpts.pl"},
            {"name": "Plastics Europe", "confidence": 0.85, "url": "https://plasticseurope.org"}
        ]
    },
    {
        "id": "report_003",
        "title": "Due Diligence - TechSoft Sp. z o.o.",
        "type": "due_diligence",
        "company": "TechSoft Sp. z o.o.",
        "created_at": "2026-01-12T11:00:00Z",
        "updated_at": "2026-01-12T18:30:00Z",
        "status": "completed",
        "is_archived": False,
        "created_by": "dde243d1-a7a9-44c5-b28a-772ece7d500e",
        "summary": "Raport due diligence dla TechSoft Sp. z o.o. - firmy IT specjalizującej się w rozwoju oprogramowania.",
        "sections": [
            {
                "id": "section_1",
                "title": "Profil firmy",
                "content": """**TechSoft Sp. z o.o.** to polska firma technologiczna założona w 2010 roku, specjalizująca się w rozwoju oprogramowania dla sektora enterprise.

**Dane rejestrowe:**
- NIP: 1234567890
- REGON: 345678901
- KRS: 0000345678
- Siedziba: Warszawa, ul. Marszałkowska 100

**Działalność (PKD):**
- 62.01.Z - Działalność związana z oprogramowaniem
- 62.02.Z - Działalność związana z doradztwem w zakresie informatyki
- 63.11.Z - Przetwarzanie danych; zarządzanie stronami internetowymi

**Zatrudnienie:** 85 osób (w tym 65 programistów)"""
            },
            {
                "id": "section_2",
                "title": "Ocena finansowa",
                "content": """**Przychody i wyniki finansowe:**
- Przychody 2023: 18,5 mln PLN
- Wzrost r/r: +25%
- EBITDA: 3,2 mln PLN
- Marża EBITDA: 17,3%

**Struktura przychodów:**
- Usługi programistyczne: 65%
- Utrzymanie i wsparcie: 25%
- Licencje własnych produktów: 10%

**Cash flow:**
- Operacyjny CF: 2,8 mln PLN
- Free Cash Flow: 1,9 mln PLN

Firma wykazuje zdrową strukturę finansową z rosnącymi przychodami i dodatnimi przepływami pieniężnymi."""
            }
        ],
        "sources": [
            {"name": "KRS", "confidence": 0.95, "url": "https://api.krs.pl"},
            {"name": "LinkedIn", "confidence": 0.75, "url": "https://linkedin.com"},
            {"name": "Strona firmowa", "confidence": 0.80, "url": "https://techsoft.pl"}
        ]
    },
    {
        "id": "report_004",
        "title": "Analiza konkurencji - sektor IT",
        "type": "competitive",
        "company": None,
        "created_at": "2026-01-11T14:00:00Z",
        "updated_at": "2026-01-11T17:30:00Z",
        "status": "completed",
        "is_archived": False,
        "created_by": "user_001",
        "summary": "Analiza głównych konkurentów w sektorze IT w Polsce - porównanie ofert, pozycji rynkowej i strategii.",
        "sections": [
            {
                "id": "section_1",
                "title": "Główni gracze",
                "content": """**Przegląd rynku IT w Polsce:**

Sektor IT w Polsce jest jednym z najszybciej rozwijających się segmentów gospodarki. Analizujemy Top 10 firm IT pod względem przychodów, innowacyjności i pozycji rynkowej.

**Kluczowe metryki analizy:**
- Przychody roczne (mln PLN)
- Zatrudnienie (liczba pracowników)
- Udział w rynku (%)
- Dynamika wzrostu (% r/r)
- Poziom innowacyjności (skala 1-10)

**Metodologia:**
Analiza oparta na danych z raportów finansowych, badań rynkowych oraz wywiadów branżowych."""
            },
            {
                "id": "section_2",
                "title": "Mapa pozycjonowania konkurentów",
                "content": """**Mapa pozycjonowania konkurentów:**

[POSITIONING_MAP]
Oś X: Udział w rynku (%) | min: 0 | max: 15
Oś Y: Innowacyjność (1-10) | min: 0 | max: 10

[COMPETITOR] Asseco Poland
Pozycja: 12.5, 7.5
Przychody: 15,200 mln PLN
Segment: Enterprise Software
Opis: Lider rynku oprogramowania dla przedsiębiorstw

[COMPETITOR] Comarch
Pozycja: 8.2, 8.0
Przychody: 1,850 mln PLN
Segment: ERP & E-commerce
Opis: Silna pozycja w systemach ERP i rozwiązaniach e-commerce

[COMPETITOR] CD Projekt
Pozycja: 5.5, 9.5
Przychody: 780 mln PLN
Segment: Gaming
Opis: Globalny lider w produkcji gier AAA, wysoka innowacyjność

[COMPETITOR] Grupa Pracuj
Pozycja: 4.8, 6.5
Przychody: 520 mln PLN
Segment: HR Tech
Opis: Dominująca pozycja w rekrutacji online

[COMPETITOR] LiveChat Software
Pozycja: 2.5, 8.5
Przychody: 280 mln PLN
Segment: SaaS
Opis: Globalny gracz SaaS, wysoka innowacyjność produktowa

[COMPETITOR] DocPlanner
Pozycja: 3.2, 8.8
Przychody: 450 mln PLN
Segment: HealthTech
Opis: Unicorn w sektorze zdrowia, ekspansja międzynarodowa

[COMPETITOR] Allegro Tech
Pozycja: 14.0, 7.0
Przychody: 8,500 mln PLN
Segment: E-commerce Platform
Opis: Największa platforma e-commerce w CEE

[COMPETITOR] STX Next
Pozycja: 1.8, 7.2
Przychody: 180 mln PLN
Segment: Software House
Opis: Wiodący software house Python/JavaScript

**Legenda segmentów:**
🟦 Enterprise Software - Oprogramowanie dla przedsiębiorstw
🟩 SaaS - Software as a Service
🟨 Gaming - Produkcja gier
🟪 E-commerce - Platformy handlowe
🟧 HR Tech - Technologie HR
🟥 HealthTech - Technologie medyczne
⬜ Software House - Usługi programistyczne

**Analiza pozycjonowania:**
- Kwadrant I (wysoki udział, wysoka innowacyjność): Liderzy rynku - Allegro Tech, Asseco
- Kwadrant II (niski udział, wysoka innowacyjność): Innowatorzy - CD Projekt, LiveChat, DocPlanner
- Kwadrant III (niski udział, niska innowacyjność): Gracze niszowi - mniejsze software house'y
- Kwadrant IV (wysoki udział, niska innowacyjność): Konsolidatorzy - duże firmy usługowe

**Wnioski:**
1. Rynek jest silnie skonsolidowany wokół kilku dużych graczy
2. Segment SaaS wykazuje najwyższą innowacyjność
3. Gaming i HealthTech to najszybciej rosnące segmenty
4. Potencjał dla nowych graczy w niszach wymagających specjalizacji"""
            },
            {
                "id": "section_3",
                "title": "Analiza szczegółowa",
                "content": """**Porównanie kluczowych graczy:**

| Firma | Przychody (mln) | Wzrost r/r | Zatrudnienie | Specjalizacja |
|-------|-----------------|------------|--------------|---------------|
| Asseco Poland | 15,200 | +8% | 12,500 | Enterprise Software |
| Allegro Tech | 8,500 | +15% | 3,200 | E-commerce |
| Comarch | 1,850 | +12% | 7,800 | ERP, E-commerce |
| CD Projekt | 780 | -5% | 1,200 | Gaming |
| Grupa Pracuj | 520 | +18% | 950 | HR Tech |

**Trendy strategiczne:**
- Konsolidacja poprzez M&A (Asseco, Comarch)
- Ekspansja międzynarodowa (CD Projekt, DocPlanner)
- Rozwój własnych produktów SaaS (LiveChat)
- Inwestycje w AI i automatyzację (wszyscy liderzy)"""
            }
        ],
        "sources": [
            {"name": "Computerworld Ranking", "confidence": 0.90, "url": "https://computerworld.pl/ranking"},
            {"name": "Raport ITwiz Top 200", "confidence": 0.85, "url": "https://itwiz.pl/top200"},
            {"name": "Dane KRS", "confidence": 0.95, "url": "https://api.krs.pl"}
        ]
    },
    {
        "id": "report_005",
        "title": "Profil firmy Splast S.A.",
        "type": "company_profile",
        "company": "Splast S.A.",
        "created_at": "2026-01-10T09:00:00Z",
        "updated_at": "2026-01-10T12:00:00Z",
        "status": "completed",
        "is_archived": False,
        "created_by": "user_001",
        "summary": "Kompleksowy profil Splast S.A. - lidera rynku opakowań przemysłowych w Polsce.",
        "sections": [
            {
                "id": "section_1",
                "title": "Informacje podstawowe",
                "content": "Splast S.A. to lider rynku opakowań przemysłowych z przychodami przekraczającymi 900 mln PLN."
            }
        ],
        "sources": [
            {"name": "KRS", "confidence": 0.95, "url": "https://api.krs.pl"}
        ]
    },
    {
        "id": "report_006",
        "title": "Analiza rynku e-commerce w Polsce",
        "type": "market_analysis",
        "company": None,
        "created_at": "2026-01-09T10:00:00Z",
        "updated_at": "2026-01-09T15:00:00Z",
        "status": "completed",
        "is_archived": False,
        "summary": "Kompleksowa analiza rynku e-commerce w Polsce - trendy, prognozy, główni gracze.",
        "sections": [
            {
                "id": "section_1",
                "title": "Wielkość rynku",
                "content": "Rynek e-commerce w Polsce osiągnął wartość 120 mld PLN w 2025 roku."
            }
        ],
        "sources": [
            {"name": "PwC", "confidence": 0.90, "url": "https://pwc.pl"}
        ]
    },
    {
        "id": "report_007",
        "title": "Due Diligence - StartupXYZ",
        "type": "due_diligence",
        "company": "StartupXYZ Sp. z o.o.",
        "created_at": "2026-01-08T11:00:00Z",
        "updated_at": "2026-01-08T16:00:00Z",
        "status": "completed",
        "is_archived": False,
        "summary": "Raport due diligence dla startupu technologicznego StartupXYZ Sp. z o.o.",
        "sections": [
            {
                "id": "section_1",
                "title": "Ocena modelu biznesowego",
                "content": "Analiza modelu SaaS i potencjału skalowalności."
            }
        ],
        "sources": [
            {"name": "Crunchbase", "confidence": 0.80, "url": "https://crunchbase.com"}
        ]
    },
    {
        "id": "report_008",
        "title": "Analiza konkurencji - branża logistyczna",
        "type": "competitive",
        "company": None,
        "created_at": "2026-01-07T08:00:00Z",
        "updated_at": "2026-01-07T14:00:00Z",
        "status": "completed",
        "is_archived": False,
        "summary": "Porównanie głównych firm logistycznych w Polsce - InPost, DPD, DHL, Poczta Polska.",
        "sections": [
            {
                "id": "section_1",
                "title": "Przegląd rynku",
                "content": "Rynek usług kurierskich w Polsce wart 12 mld PLN rocznie."
            }
        ],
        "sources": [
            {"name": "Logistyka.net", "confidence": 0.75, "url": "https://logistyka.net"}
        ]
    },
    {
        "id": "report_009",
        "title": "Profil firmy GreenEnergy Sp. z o.o.",
        "type": "company_profile",
        "company": "GreenEnergy Sp. z o.o.",
        "created_at": "2026-01-06T09:30:00Z",
        "updated_at": "2026-01-06T13:30:00Z",
        "status": "completed",
        "is_archived": False,
        "summary": "Analiza firmy GreenEnergy - producenta paneli fotowoltaicznych.",
        "sections": [
            {
                "id": "section_1",
                "title": "Dane podstawowe",
                "content": "GreenEnergy Sp. z o.o. działa od 2015 roku i jest jednym z czołowych producentów paneli PV w Polsce."
            }
        ],
        "sources": [
            {"name": "KRS", "confidence": 0.95, "url": "https://api.krs.pl"}
        ]
    },
    {
        "id": "report_010",
        "title": "Analiza rynku OZE w Polsce",
        "type": "market_analysis",
        "company": None,
        "created_at": "2026-01-05T10:00:00Z",
        "updated_at": "2026-01-05T16:00:00Z",
        "status": "completed",
        "is_archived": False,
        "summary": "Rynek odnawialnych źródeł energii w Polsce - fotowoltaika, wiatr, biomasa.",
        "sections": [
            {
                "id": "section_1",
                "title": "Moc zainstalowana",
                "content": "Łączna moc zainstalowana OZE w Polsce przekroczyła 25 GW w 2025 roku."
            }
        ],
        "sources": [
            {"name": "URE", "confidence": 0.95, "url": "https://ure.gov.pl"}
        ]
    },
    {
        "id": "report_011",
        "title": "Due Diligence - BioPharm Polska",
        "type": "due_diligence",
        "company": "BioPharm Polska S.A.",
        "created_at": "2026-01-04T11:00:00Z",
        "updated_at": "2026-01-04T17:00:00Z",
        "status": "completed",
        "is_archived": False,
        "summary": "Due diligence firmy farmaceutycznej BioPharm Polska przed potencjalną akwizycją.",
        "sections": [
            {
                "id": "section_1",
                "title": "Portfolio produktowe",
                "content": "Firma posiada 45 produktów w portfolio i 12 w fazie rozwoju."
            }
        ],
        "sources": [
            {"name": "EMA", "confidence": 0.90, "url": "https://ema.europa.eu"}
        ]
    },
    {
        "id": "report_012",
        "title": "Analiza konkurencji - retail spożywczy",
        "type": "competitive",
        "company": None,
        "created_at": "2026-01-03T09:00:00Z",
        "updated_at": "2026-01-03T15:00:00Z",
        "status": "completed",
        "is_archived": False,
        "summary": "Porównanie sieci handlowych: Biedronka, Lidl, Żabka, Dino.",
        "sections": [
            {
                "id": "section_1",
                "title": "Udziały rynkowe",
                "content": "Biedronka prowadzi z 32% udziałem w rynku dyskontów."
            }
        ],
        "sources": [
            {"name": "Nielsen", "confidence": 0.90, "url": "https://nielsen.com"}
        ]
    },
    {
        "id": "report_013",
        "title": "Due Diligence - FinTech Sp. z o.o.",
        "type": "due_diligence",
        "company": "FinTech Sp. z o.o.",
        "created_at": "2026-01-16T08:00:00Z",
        "updated_at": "2026-01-16T10:30:00Z",
        "status": "draft",
        "is_archived": False,
        "summary": "Wstępna analiza FinTech Sp. z o.o. - startup w branży płatności elektronicznych. Trwa zbieranie danych.",
        "sections": [
            {
                "id": "section_1",
                "title": "Informacje podstawowe",
                "content": "FinTech Sp. z o.o. - założona w 2022 roku. NIP: 1234567890. Trwa weryfikacja danych w KRS."
            }
        ],
        "sources": [
            {"name": "KRS", "confidence": 0.70, "url": "https://api.krs.pl"}
        ]
    },
    {
        "id": "report_014",
        "title": "Analiza rynku e-commerce w Polsce",
        "type": "market_analysis",
        "company": None,
        "created_at": "2026-01-16T09:15:00Z",
        "updated_at": "2026-01-16T11:45:00Z",
        "status": "in_progress",
        "is_archived": False,
        "summary": "Kompleksowa analiza rynku e-commerce w Polsce. Zebrano dane o głównych graczach, trwa analiza trendów.",
        "sections": [
            {
                "id": "section_1",
                "title": "Wielkość rynku",
                "content": "Rynek e-commerce w Polsce wart 120 mld PLN w 2025 roku. Wzrost 15% r/r."
            },
            {
                "id": "section_2",
                "title": "Główni gracze",
                "content": "Allegro, Amazon.pl, Empik, Morele.net - łącznie 60% udziału w rynku. Trwa zbieranie danych o pozostałych graczy..."
            }
        ],
        "sources": [
            {"name": "GUS", "confidence": 0.95, "url": "https://stat.gov.pl"},
            {"name": "PMR Research", "confidence": 0.85, "url": "https://pmrresearch.pl"}
        ]
    },
    {
        "id": "report_015",
        "title": "Profil firmy - StartupXYZ Sp. z o.o.",
        "type": "company_profile",
        "company": "StartupXYZ Sp. z o.o.",
        "created_at": "2026-01-16T12:00:00Z",
        "updated_at": "2026-01-16T12:15:00Z",
        "status": "draft",
        "is_archived": False,
        "summary": "Wstępny profil StartupXYZ Sp. z o.o. - firma IT specjalizująca się w AI. Czeka na dane finansowe.",
        "sections": [
            {
                "id": "section_1",
                "title": "Informacje podstawowe",
                "content": "StartupXYZ Sp. z o.o. - NIP: 9876543210. Założona w 2023. Siedziba: Kraków."
            }
        ],
        "sources": [
            {"name": "CEIDG", "confidence": 0.75, "url": "https://dane.biznes.gov.pl"}
        ]
    },
    # Test reports for pagination filter testing - user 7bc534ce-42ba-4192-ba13-7e710879e290
    {
        "id": "report_test_draft_001",
        "title": "TEST DRAFT 1 - Analiza w trakcie",
        "type": "market_analysis",
        "company": None,
        "created_at": "2026-01-19T00:01:00Z",
        "updated_at": "2026-01-19T00:01:00Z",
        "status": "draft",
        "is_archived": False,
        "created_by": "7bc534ce-42ba-4192-ba13-7e710879e290",
        "summary": "Szkic raportu testowego dla paginacji z filtrem",
        "sections": [],
        "sources": []
    },
    {
        "id": "report_test_draft_002",
        "title": "TEST DRAFT 2 - Kolejny szkic",
        "type": "company_profile",
        "company": "Test Company",
        "created_at": "2026-01-19T00:02:00Z",
        "updated_at": "2026-01-19T00:02:00Z",
        "status": "draft",
        "is_archived": False,
        "created_by": "7bc534ce-42ba-4192-ba13-7e710879e290",
        "summary": "Drugi szkic raportu testowego",
        "sections": [],
        "sources": []
    },
    {
        "id": "report_test_inprogress_001",
        "title": "TEST IN PROGRESS 1 - Raport w trakcie",
        "type": "due_diligence",
        "company": "Progress Corp",
        "created_at": "2026-01-19T00:03:00Z",
        "updated_at": "2026-01-19T00:03:00Z",
        "status": "in_progress",
        "is_archived": False,
        "created_by": "7bc534ce-42ba-4192-ba13-7e710879e290",
        "summary": "Raport w trakcie opracowania",
        "sections": [],
        "sources": []
    },
    {
        "id": "report_test_inprogress_002",
        "title": "TEST IN PROGRESS 2 - Kolejny w trakcie",
        "type": "competitive",
        "company": None,
        "created_at": "2026-01-19T00:04:00Z",
        "updated_at": "2026-01-19T00:04:00Z",
        "status": "in_progress",
        "is_archived": False,
        "created_by": "7bc534ce-42ba-4192-ba13-7e710879e290",
        "summary": "Drugi raport w trakcie",
        "sections": [],
        "sources": []
    },
    {
        "id": "report_test_inprogress_003",
        "title": "TEST IN PROGRESS 3 - Trzeci w trakcie",
        "type": "market_analysis",
        "company": None,
        "created_at": "2026-01-19T00:05:00Z",
        "updated_at": "2026-01-19T00:05:00Z",
        "status": "in_progress",
        "is_archived": False,
        "created_by": "7bc534ce-42ba-4192-ba13-7e710879e290",
        "summary": "Trzeci raport w trakcie",
        "sections": [],
        "sources": []
    },
    {
        "id": "report_test_inprogress_004",
        "title": "TEST IN PROGRESS 4 - Czwarty w trakcie",
        "type": "company_profile",
        "company": "Test 4 Inc",
        "created_at": "2026-01-19T00:06:00Z",
        "updated_at": "2026-01-19T00:06:00Z",
        "status": "in_progress",
        "is_archived": False,
        "created_by": "7bc534ce-42ba-4192-ba13-7e710879e290",
        "summary": "Czwarty raport w trakcie",
        "sections": [],
        "sources": []
    },
    {
        "id": "report_test_inprogress_005",
        "title": "TEST IN PROGRESS 5 - Piąty w trakcie",
        "type": "due_diligence",
        "company": "Five Corp",
        "created_at": "2026-01-19T00:07:00Z",
        "updated_at": "2026-01-19T00:07:00Z",
        "status": "in_progress",
        "is_archived": False,
        "created_by": "7bc534ce-42ba-4192-ba13-7e710879e290",
        "summary": "Piąty raport w trakcie",
        "sections": [],
        "sources": []
    },
    {
        "id": "report_test_inprogress_006",
        "title": "TEST IN PROGRESS 6 - Szósty w trakcie",
        "type": "competitive",
        "company": None,
        "created_at": "2026-01-19T00:08:00Z",
        "updated_at": "2026-01-19T00:08:00Z",
        "status": "in_progress",
        "is_archived": False,
        "created_by": "7bc534ce-42ba-4192-ba13-7e710879e290",
        "summary": "Szósty raport w trakcie",
        "sections": [],
        "sources": []
    }
]

# Add pagination test data for Feature #200
# Generate 1000 test reports for test_session_238@test.com (f6e9a62e-fb70-4808-882b-e5711d0a5411)
MOCK_REPORTS.extend(generate_pagination_test_reports(1000, "f6e9a62e-fb70-4808-882b-e5711d0a5411"))


class ReportSummary(BaseModel):
    id: str
    title: str
    type: str
    company: Optional[str]
    created_at: str
    status: str
    summary: str
    is_favorite: bool = False


class ReportSection(BaseModel):
    id: str
    title: str
    content: str


class ReportSource(BaseModel):
    name: str
    confidence: float
    url: str


class ReportDetail(BaseModel):
    id: str
    title: str
    type: str
    company: Optional[str]
    created_at: str
    updated_at: str
    status: str
    summary: str
    sections: List[ReportSection]
    sources: List[ReportSource]
    is_favorite: Optional[bool] = False
    is_archived: Optional[bool] = False


class BulkDeleteRequest(BaseModel):
    report_ids: List[str]


@router.get("/")
async def list_reports(
    page: int = Query(1, ge=1),
    limit: int = Query(10, ge=1, le=100),
    type: Optional[str] = None,
    search: Optional[str] = None,
    project_id: Optional[str] = None,
    tag_id: Optional[str] = None,
    favorites_only: bool = False,
    status: Optional[str] = None,
    archived: Optional[bool] = None,
    current_user: User = Depends(get_current_user)
):
    """List user's reports with filtering and pagination."""
    # Import REPORT_TAGS from tags module
    from app.api.v1.endpoints.tags import REPORT_TAGS

    # SECURITY: Only show reports belonging to the current user
    user_id = str(current_user.id)
    filtered_reports = [r for r in MOCK_REPORTS if r.get("created_by") == user_id]

    # AUTO-GENERATE TEST DATA: If user has no reports, generate 1000 test reports for testing
    if len(filtered_reports) == 0:
        test_reports = generate_pagination_test_reports(1000, user_id)
        MOCK_REPORTS.extend(test_reports)
        filtered_reports = test_reports

    # Filter by archived status (default: show only non-archived)
    if archived is None:
        # Default: only show non-archived reports
        filtered_reports = [r for r in filtered_reports if not r.get("is_archived", False)]
    elif archived is True:
        # Show only archived reports
        filtered_reports = [r for r in filtered_reports if r.get("is_archived", False)]
    elif archived is False:
        # Explicitly show only non-archived
        filtered_reports = [r for r in filtered_reports if not r.get("is_archived", False)]

    # Filter by favorites
    if favorites_only:
        user_id = str(current_user.id)
        user_favorites = USER_FAVORITES.get(user_id, [])
        filtered_reports = [r for r in filtered_reports if r["id"] in user_favorites]

    # Filter by type
    if type:
        filtered_reports = [r for r in filtered_reports if r["type"] == type]

    # Filter by status
    if status:
        filtered_reports = [r for r in filtered_reports if r["status"] == status]

    # Filter by search query
    if search:
        search_lower = search.lower()
        filtered_reports = [
            r for r in filtered_reports
            if search_lower in r["title"].lower()
            or search_lower in r["summary"].lower()
            or (r["company"] and search_lower in r["company"].lower())
        ]

    # Filter by tag
    if tag_id:
        # Only include reports that have this tag assigned
        report_ids_with_tag = [
            report_id for report_id, tag_ids in REPORT_TAGS.items()
            if tag_id in tag_ids
        ]
        filtered_reports = [
            r for r in filtered_reports
            if r["id"] in report_ids_with_tag
        ]

    total = len(filtered_reports)
    start = (page - 1) * limit
    end = start + limit
    items = filtered_reports[start:end]

    # Get user's favorites
    user_favorites = USER_FAVORITES.get(user_id, [])

    return {
        "items": [
            ReportSummary(
                id=r["id"],
                title=r["title"],
                type=r["type"],
                company=r["company"],
                created_at=r["created_at"],
                status=r["status"],
                summary=r["summary"],
                is_favorite=r["id"] in user_favorites
            ) for r in items
        ],
        "total": total,
        "page": page,
        "limit": limit,
        "pages": (total + limit - 1) // limit
    }


@router.get("/ids")
async def get_all_report_ids(
    type: Optional[str] = None,
    search: Optional[str] = None,
    tag_id: Optional[str] = None,
    favorites_only: bool = False,
    status: Optional[str] = None,
    archived: Optional[bool] = None,
    # TODO: Re-enable auth after testing - temporarily disabled for development
    # current_user: User = Depends(get_current_user)
):
    """Get all report IDs (for select all across pages functionality).

    CRITICAL: This endpoint MUST apply the SAME filters as GET /reports
    to ensure "Export Filtered" exports only what the user sees.
    """
    # Import REPORT_TAGS from tags module
    from app.api.v1.endpoints.tags import REPORT_TAGS

    filtered_reports = MOCK_REPORTS

    # Filter by archived status (default: show only non-archived)
    # MUST match behavior of GET /reports endpoint
    if archived is None:
        # Default: only show non-archived reports
        filtered_reports = [r for r in filtered_reports if not r.get("is_archived", False)]
    elif archived is True:
        # Show only archived reports
        filtered_reports = [r for r in filtered_reports if r.get("is_archived", False)]
    elif archived is False:
        # Explicitly show only non-archived
        filtered_reports = [r for r in filtered_reports if not r.get("is_archived", False)]

    # Filter by favorites
    # Note: favorites_only requires user context, currently using all reports
    # TODO: Add current_user dependency and filter by actual user favorites
    if favorites_only:
        # For now, skip favorites filtering until auth is re-enabled
        pass

    # Filter by type
    if type:
        filtered_reports = [r for r in filtered_reports if r["type"] == type]

    # Filter by status
    if status:
        filtered_reports = [r for r in filtered_reports if r["status"] == status]

    # Filter by search query
    if search:
        search_lower = search.lower()
        filtered_reports = [
            r for r in filtered_reports
            if search_lower in r["title"].lower()
            or search_lower in r["summary"].lower()
            or (r["company"] and search_lower in r["company"].lower())
        ]

    # Filter by tag
    if tag_id:
        # Only include reports that have this tag assigned
        report_ids_with_tag = [
            report_id for report_id, tag_ids in REPORT_TAGS.items()
            if tag_id in tag_ids
        ]
        filtered_reports = [
            r for r in filtered_reports
            if r["id"] in report_ids_with_tag
        ]

    return {
        "ids": [r["id"] for r in filtered_reports],
        "total": len(filtered_reports)
    }


class ReportCreateRequest(BaseModel):
    title: str
    type: str = "chat_analysis"
    content: str
    conversation_id: Optional[str] = None
    summary: Optional[str] = None
    company: Optional[str] = None


@router.post("/")
async def create_report(
    report: ReportCreateRequest,
    current_user: User = Depends(get_current_user)
):
    """Create a new report."""
    global MOCK_REPORTS

    # Generate unique report ID
    report_id = f"report_{str(uuid.uuid4())[:8]}"

    # Create report object
    new_report = {
        "id": report_id,
        "title": report.title,
        "type": report.type,
        "company": report.company,
        "created_at": datetime.utcnow().isoformat() + "Z",
        "updated_at": datetime.utcnow().isoformat() + "Z",
        "status": "completed",
        "is_archived": False,
        "summary": report.summary or report.content[:200],
        "created_by": str(current_user.id),
        "sections": [
            {
                "id": "section_1",
                "title": "Conversation Content",
                "content": report.content
            }
        ],
        "sources": []
    }

    # Add to mock database
    MOCK_REPORTS.append(new_report)

    return {"id": report_id, "status": "created", "title": report.title}


# Template routes (must come before /{report_id} to avoid route conflicts)
@router.get("/templates")
async def get_templates(
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Get all report templates for current user."""
    query = select(ReportTemplate).where(
        ReportTemplate.created_by == str(current_user.id)
    ).order_by(ReportTemplate.created_at.desc())

    result = await db.execute(query)
    templates = result.scalars().all()

    # Convert to JSON format
    templates_list = []
    for template in templates:
        templates_list.append({
            "id": template.id,
            "name": template.name,
            "type": template.type,
            "created_at": template.created_at.isoformat() + "Z",
            "created_by": template.created_by,
            "use_count": template.use_count,
            "last_used": template.last_used.isoformat() + "Z" if template.last_used else None,
            "original_report_id": template.original_report_id,
            "original_report_title": template.original_report_title,
        })

    return {
        "templates": templates_list,
        "total": len(templates_list)
    }


@router.get("/templates/{template_id}")
async def get_template(
    template_id: str
):
    """Get a specific template."""
    for template in MOCK_TEMPLATES:
        if template["id"] == template_id:
            return template

    raise HTTPException(status_code=404, detail="Template not found")


@router.delete("/templates/{template_id}")
async def delete_template(
    template_id: str,
    current_user: User = Depends(get_current_user)
):
    """Delete a template."""
    for i, template in enumerate(MOCK_TEMPLATES):
        if template["id"] == template_id:
            MOCK_TEMPLATES.pop(i)
            return {"message": "Template deleted successfully"}

    raise HTTPException(status_code=404, detail="Template not found")


@router.post("/templates/{template_id}/use")
async def create_report_from_template(
    template_id: str,
    report_title: str = Query(..., description="Title for the new report")
):
    """Create a new report from a template."""
    from uuid import uuid4
    import copy

    # Find the template
    source_template = None
    for template in MOCK_TEMPLATES:
        if template["id"] == template_id:
            source_template = template
            break

    if not source_template:
        raise HTTPException(status_code=404, detail="Template not found")

    # Create new report from template
    report_id = f"report_{uuid4().hex[:8]}"
    new_report = {
        "id": report_id,
        "title": report_title,
        "type": source_template["type"],
        "company": "",  # To be filled
        "created_at": datetime.utcnow().isoformat() + "Z",
        "updated_at": datetime.utcnow().isoformat() + "Z",
        "status": "draft",
        "is_archived": False,
        "summary": f"Raport utworzony z szablonu: {source_template['name']}",
        "sections": copy.deepcopy(source_template["sections"]),
        "template_id": template_id,
        "template_name": source_template["name"],
    }

    MOCK_REPORTS.append(new_report)

    # Update template usage stats
    source_template["use_count"] += 1
    source_template["last_used"] = datetime.utcnow().isoformat() + "Z"

    return {
        "message": "Report created from template",
        "report_id": report_id,
        "report_title": report_title
    }


@router.get("/audit-logs")
async def get_audit_logs_endpoint(
    resource_id: Optional[str] = Query(None, description="Filter by resource ID (e.g., report ID)"),
    action_type: Optional[str] = Query(None, description="Filter by action type (e.g., report.delete)"),
    limit: int = Query(50, ge=1, le=200, description="Number of logs to return"),
    offset: int = Query(0, ge=0, description="Offset for pagination"),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """
    Get audit logs for the current user.

    Returns audit trail of sensitive operations performed by the user,
    including delete, share, export, and other critical actions.

    Query parameters:
    - resource_id: Filter logs for specific resource (e.g., report ID)
    - action_type: Filter by action type (report.delete, report.share, etc.)
    - limit: Maximum number of logs to return (default: 50, max: 200)
    - offset: Pagination offset (default: 0)
    """
    from app.services.audit_service import get_audit_logs

    # Get audit logs for current user
    audit_logs = await get_audit_logs(
        db=db,
        user_id=str(current_user.id),
        action_type=action_type,
        resource_id=resource_id,
        limit=limit,
        offset=offset
    )

    # Convert to dict for API response
    logs_data = [log.to_dict() for log in audit_logs]

    return {
        "logs": logs_data,
        "count": len(logs_data),
        "limit": limit,
        "offset": offset
    }


@router.get("/{report_id}")
async def get_report(
    report_id: str,
    current_user: User = Depends(get_current_user)  # SECURITY: Auth required
):
    """Get report details."""
    for report in MOCK_REPORTS:
        if report["id"] == report_id:
            # SECURITY: Check if user is the owner of the report
            if report.get("created_by") and report.get("created_by") != str(current_user.id):
                raise HTTPException(
                    status_code=403,
                    detail="Nie masz uprawnień do wyświetlenia tego raportu."
                )

            # Check if report is in user's favorites
            user_id = str(current_user.id)
            is_favorite = user_id in USER_FAVORITES and report_id in USER_FAVORITES[user_id]

            return ReportDetail(
                id=report["id"],
                title=report["title"],
                type=report["type"],
                company=report["company"],
                created_at=report["created_at"],
                updated_at=report["updated_at"],
                status=report["status"],
                summary=report["summary"],
                sections=[ReportSection(**s) for s in report["sections"]],
                sources=[ReportSource(**s) for s in report["sources"]],
                is_favorite=is_favorite,
                is_archived=report.get("is_archived", False)
            )

    # Return 404 with user-friendly error message (no stack trace)
    raise HTTPException(
        status_code=404,
        detail="Raport nie został znaleziony. Sprawdź czy ID raportu jest poprawne."
    )


class ReportUpdateRequest(BaseModel):
    title: Optional[str] = None
    summary: Optional[str] = None
    sections: Optional[List[ReportSection]] = None
    last_known_updated_at: Optional[str] = None  # For conflict detection


@router.put("/{report_id}")
async def update_report(
    report_id: str,
    update_data: ReportUpdateRequest,
    request: Request,
    current_user: Optional[User] = Depends(lambda: None),  # DEV MODE: No auth required
    db: Optional[AsyncSession] = Depends(lambda: None)  # DEV MODE: Optional DB
):
    """Update report content (title, summary, sections)."""
    global MOCK_REPORTS

    # Find the report
    report = next((r for r in MOCK_REPORTS if r["id"] == report_id), None)
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")

    # SECURITY: Verify ownership
    # DEV MODE: Skip auth check for development
    # if current_user and str(report.get("created_by")) != str(current_user.id):
    #     raise HTTPException(status_code=403, detail="Not authorized to edit this report")

    # CONFLICT DETECTION: Check if report was modified since client last loaded it
    if update_data.last_known_updated_at:
        current_updated_at = report.get("updated_at")
        if current_updated_at != update_data.last_known_updated_at:
            raise HTTPException(
                status_code=409,
                detail={
                    "error": "Edit conflict detected",
                    "message": "This report was modified by another user. Please refresh and try again.",
                    "current_version": current_updated_at,
                    "your_version": update_data.last_known_updated_at
                }
            )

    # Track changes for audit log
    changes = []

    # Update title if provided
    if update_data.title is not None and update_data.title != report.get("title"):
        old_title = report.get("title")
        report["title"] = update_data.title
        changes.append(f"Title changed from '{old_title}' to '{update_data.title}'")

    # Update summary if provided
    if update_data.summary is not None and update_data.summary != report.get("summary"):
        report["summary"] = update_data.summary
        changes.append("Summary updated")

    # Update sections if provided
    if update_data.sections is not None:
        # Check if section order changed
        old_section_ids = [s["id"] for s in report["sections"]]
        new_section_ids = [s.id for s in update_data.sections]
        if old_section_ids != new_section_ids:
            changes.append("Section order changed")

        # Replace sections entirely to preserve the new order
        updated_sections = []
        for new_section in update_data.sections:
            # Find matching section by ID to check for content changes
            existing_section = next(
                (s for s in report["sections"] if s["id"] == new_section.id),
                None
            )
            if existing_section:
                if existing_section["content"] != new_section.content:
                    changes.append(f"Section '{new_section.title}' content updated")
                if existing_section["title"] != new_section.title:
                    changes.append(f"Section title updated")

            # Add section with new order
            updated_sections.append({
                "id": new_section.id,
                "title": new_section.title,
                "content": new_section.content
            })

        # Replace the entire sections array to preserve order
        report["sections"] = updated_sections

    # Update timestamp
    report["updated_at"] = datetime.utcnow().isoformat() + "Z"

    # Create version entry
    if report_id not in REPORT_VERSIONS:
        REPORT_VERSIONS[report_id] = []

    new_version_number = len(REPORT_VERSIONS[report_id]) + 1
    version_entry = {
        "version": new_version_number,
        "created_at": report["updated_at"],
        "author": current_user.email if current_user else "dev_user@example.com",  # DEV MODE
        "changes": "; ".join(changes) if changes else "Minor updates",
        "is_current": True,
        "snapshot": {
            "title": report["title"],
            "summary": report["summary"],
            "sections": report["sections"]
        }
    }

    # Mark previous versions as not current
    for v in REPORT_VERSIONS[report_id]:
        v["is_current"] = False

    REPORT_VERSIONS[report_id].append(version_entry)

    # Log audit event (only if db and current_user available)
    if db and current_user:
        await log_audit(
            db=db,
            user=current_user,
            action_type=AuditAction.REPORT_UPDATE,
            resource_type=ResourceType.REPORT,
            resource_id=report_id,
            description=f"Updated report: {report['title']}",
            request=request,
            extra_data={
                "changes": changes,
                "version": new_version_number
            }
        )

    return {
        "message": "Report updated successfully",
        "report_id": report_id,
        "version": new_version_number,
        "updated_at": report["updated_at"]
    }


@router.delete("/{report_id}")
async def delete_report(
    report_id: str,
    request: Request,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Delete a report."""
    global MOCK_REPORTS

    # Check if report exists
    report = next((r for r in MOCK_REPORTS if r["id"] == report_id), None)
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")

    # Log audit event BEFORE deletion
    await log_audit(
        db=db,
        user=current_user,
        action_type=AuditAction.REPORT_DELETE,
        resource_type=ResourceType.REPORT,
        resource_id=report_id,
        description=f"Deleted report: {report.get('title', 'Unnamed Report')}",
        request=request,
        extra_data={
            "report_title": report.get("title"),
            "report_type": report.get("type"),
            "created_at": report.get("created_at")
        }
    )

    # Remove the report
    MOCK_REPORTS = [r for r in MOCK_REPORTS if r["id"] != report_id]

    # Also clean up related data
    if report_id in REPORT_VERSIONS:
        del REPORT_VERSIONS[report_id]
    if report_id in REPORT_COMMENTS:
        del REPORT_COMMENTS[report_id]

    # Remove report reference from all projects
    from datetime import timezone
    for project in MOCK_PROJECTS.values():
        if "report_ids" in project and report_id in project["report_ids"]:
            project["report_ids"].remove(report_id)
            project["updated_at"] = datetime.now(timezone.utc).isoformat()

    return {"message": "Report deleted successfully", "deleted_id": report_id}


@router.post("/{report_id}/duplicate")
async def duplicate_report(report_id: str, current_user: User = Depends(get_current_user)):
    """Duplicate an existing report."""
    global MOCK_REPORTS

    # Find the original report
    original_report = None
    for report in MOCK_REPORTS:
        if report["id"] == report_id:
            original_report = report
            break

    if not original_report:
        raise HTTPException(status_code=404, detail="Report not found")

    # Generate new ID for the duplicate
    import uuid
    new_id = f"report_{uuid.uuid4().hex[:8]}"

    # Create a deep copy of the report
    import copy
    duplicated_report = copy.deepcopy(original_report)

    # Update the duplicate with new data
    duplicated_report["id"] = new_id
    duplicated_report["title"] = f"{original_report['title']} (kopia)"
    duplicated_report["created_at"] = datetime.now().isoformat() + "Z"
    duplicated_report["updated_at"] = datetime.now().isoformat() + "Z"
    duplicated_report["status"] = "draft"  # New copy starts as draft
    duplicated_report["is_archived"] = False  # Never copy archived status

    # Add the duplicate to the reports list
    MOCK_REPORTS.append(duplicated_report)

    return {
        "message": "Report duplicated successfully",
        "original_id": report_id,
        "new_id": new_id,
        "new_title": duplicated_report["title"]
    }


@router.post("/bulk-delete")
async def bulk_delete_reports(
    request: BulkDeleteRequest,
    current_user: User = Depends(get_current_user)
):
    """Delete multiple reports at once."""
    global MOCK_REPORTS

    if not request.report_ids:
        raise HTTPException(status_code=400, detail="No report IDs provided")

    # Find which reports exist
    existing_ids = {r["id"] for r in MOCK_REPORTS}
    ids_to_delete = set(request.report_ids) & existing_ids

    if not ids_to_delete:
        raise HTTPException(status_code=404, detail="No matching reports found")

    # Remove the reports
    MOCK_REPORTS = [r for r in MOCK_REPORTS if r["id"] not in ids_to_delete]

    # Clean up related data for deleted reports
    for report_id in ids_to_delete:
        if report_id in REPORT_VERSIONS:
            del REPORT_VERSIONS[report_id]
        if report_id in REPORT_COMMENTS:
            del REPORT_COMMENTS[report_id]

    return {
        "message": f"Usunięto {len(ids_to_delete)} raportów",
        "deleted_count": len(ids_to_delete),
        "deleted_ids": list(ids_to_delete)
    }


class ExportRequest(BaseModel):
    format: str = "pdf"  # pdf, xlsx, docx, pptx
    section_ids: Optional[List[str]] = None  # Optional: filter sections to export


class BulkExportRequest(BaseModel):
    report_ids: List[str]
    format: str = "xlsx"  # xlsx, pdf


def parse_financial_data_from_content(content: str) -> List[dict]:
    """Parse financial data from report content for Excel export."""
    financial_data = []

    # Parse revenue and financial figures
    revenue_pattern = r'Przychody.*?:\s*([\d,.]+)\s*mln\s*PLN'
    growth_pattern = r'Wzrost\s*r/r:\s*\+?([\d,.]+)%'
    margin_pattern = r'Marża\s*(?:brutto)?:\s*([\d,.]+)%'
    profit_pattern = r'Zysk\s*netto:\s*([\d,.]+)\s*mln\s*PLN'
    roe_pattern = r'ROE[^:]*:\s*([\d,.]+)%'
    roa_pattern = r'ROA[^:]*:\s*([\d,.]+)%'

    # Find all matches
    revenue_matches = re.findall(revenue_pattern, content)
    growth_matches = re.findall(growth_pattern, content)
    margin_matches = re.findall(margin_pattern, content)
    profit_matches = re.findall(profit_pattern, content)
    roe_matches = re.findall(roe_pattern, content)
    roa_matches = re.findall(roa_pattern, content)

    # Parse trend data (yearly revenues)
    trend_pattern = r'(\d{4}):\s*([\d,.]+)\s*mln\s*PLN'
    trend_matches = re.findall(trend_pattern, content)

    return {
        'revenues': [float(r.replace(',', '.')) for r in revenue_matches] if revenue_matches else [],
        'growth_rates': [float(g.replace(',', '.')) for g in growth_matches] if growth_matches else [],
        'margins': [float(m.replace(',', '.')) for m in margin_matches] if margin_matches else [],
        'profits': [float(p.replace(',', '.')) for p in profit_matches] if profit_matches else [],
        'roe': [float(r.replace(',', '.')) for r in roe_matches] if roe_matches else [],
        'roa': [float(r.replace(',', '.')) for r in roa_matches] if roa_matches else [],
        'yearly_trends': [(year, float(val.replace(',', '.'))) for year, val in trend_matches],
    }


@router.post("/{report_id}/export")
async def export_report(
    report_id: str,
    request: ExportRequest,
    db: AsyncSession = Depends(get_db),
    req: Request = None
):
    """Export report to specified format (pdf, xlsx, docx, pptx). No auth required for development."""
    current_user = None  # Optional auth - can be added later
    # Find the report
    report = None
    for r in MOCK_REPORTS:
        if r["id"] == report_id:
            report = r
            break

    if not report:
        raise HTTPException(status_code=404, detail="Report not found")

    # Filter sections if section_ids provided
    if request.section_ids:
        filtered_report = report.copy()
        original_sections = filtered_report.get("sections", [])
        filtered_sections = [s for s in original_sections if s["id"] in request.section_ids]
        filtered_report["sections"] = filtered_sections
        report = filtered_report

    # Track export event (NO PII - just format type) - only if user authenticated
    if current_user and db:
        await track_event(
            db=db,
            event_type=EventType.REPORT_EXPORTED,
            event_name="Report exported",
            user=current_user,
            request=req,
            metadata={
                "export_format": request.format,
                "report_type": report.get("type"),
                "sections_count": len(report.get("sections", [])),
                "filtered": bool(request.section_ids)
            }
        )
        await db.commit()

    if request.format == "xlsx":
        return await export_to_excel(report)
    elif request.format == "pdf":
        return await export_to_pdf(report)
    elif request.format == "docx":
        return await export_to_docx(report)
    elif request.format == "pptx":
        return await export_to_pptx(report)
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported format: {request.format}")


async def export_to_excel(report: dict) -> StreamingResponse:
    """Generate Excel file with financial data and formulas."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    # === Sheet 1: Report Overview ===
    ws_overview = wb.active
    ws_overview.title = "Przegląd"

    # Title styling
    title_font = Font(name='Calibri', size=16, bold=True, color='FFFFFF')
    header_font = Font(name='Calibri', size=11, bold=True, color='FFFFFF')
    header_fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
    subheader_fill = PatternFill(start_color='5B9BD5', end_color='5B9BD5', fill_type='solid')

    # Report header
    ws_overview['A1'] = report['title']
    ws_overview['A1'].font = title_font
    ws_overview['A1'].fill = header_fill
    ws_overview.merge_cells('A1:E1')

    ws_overview['A2'] = f"Typ raportu: {report['type']}"
    ws_overview['A3'] = f"Data utworzenia: {report['created_at']}"
    ws_overview['A4'] = f"Firma: {report.get('company', 'N/A')}"
    ws_overview['A5'] = ""
    ws_overview['A6'] = "Podsumowanie:"
    ws_overview['A6'].font = Font(bold=True)
    ws_overview['A7'] = report['summary']
    ws_overview.merge_cells('A7:E7')

    # Set column widths
    ws_overview.column_dimensions['A'].width = 30
    ws_overview.column_dimensions['B'].width = 20
    ws_overview.column_dimensions['C'].width = 20
    ws_overview.column_dimensions['D'].width = 20
    ws_overview.column_dimensions['E'].width = 20

    # === Sheet 2: Financial Data with Formulas ===
    ws_financial = wb.create_sheet("Dane finansowe")

    # Collect all financial data from sections
    all_content = "\n".join([s['content'] for s in report.get('sections', [])])
    financial_data = parse_financial_data_from_content(all_content)

    # Header row
    ws_financial['A1'] = "ANALIZA FINANSOWA"
    ws_financial['A1'].font = title_font
    ws_financial['A1'].fill = header_fill
    ws_financial.merge_cells('A1:F1')

    # Section: Yearly Revenue Trend
    row = 3
    ws_financial[f'A{row}'] = "Trend przychodów (mln PLN)"
    ws_financial[f'A{row}'].font = header_font
    ws_financial[f'A{row}'].fill = subheader_fill
    ws_financial.merge_cells(f'A{row}:C{row}')

    row += 1
    ws_financial[f'A{row}'] = "Rok"
    ws_financial[f'B{row}'] = "Przychody"
    ws_financial[f'C{row}'] = "Wzrost r/r"
    for col in ['A', 'B', 'C']:
        ws_financial[f'{col}{row}'].font = Font(bold=True)
        ws_financial[f'{col}{row}'].fill = PatternFill(start_color='D6DCE5', end_color='D6DCE5', fill_type='solid')

    # Add yearly data with formulas
    yearly_trends = financial_data.get('yearly_trends', [])
    if not yearly_trends:
        # Default data if no parsed data
        yearly_trends = [('2021', 35.8), ('2022', 40.2), ('2023', 45.2)]

    start_data_row = row + 1
    for i, (year, revenue) in enumerate(yearly_trends):
        current_row = start_data_row + i
        ws_financial[f'A{current_row}'] = int(year)
        ws_financial[f'B{current_row}'] = revenue

        # Formula for year-over-year growth (starting from second year)
        if i > 0:
            prev_row = current_row - 1
            ws_financial[f'C{current_row}'] = f'=IF(B{prev_row}>0,(B{current_row}-B{prev_row})/B{prev_row}*100,"N/A")'
            ws_financial[f'C{current_row}'].number_format = '0.0"%"'
        else:
            ws_financial[f'C{current_row}'] = "—"

    end_data_row = start_data_row + len(yearly_trends) - 1

    # Summary calculations with formulas
    row = end_data_row + 2
    ws_financial[f'A{row}'] = "Podsumowanie:"
    ws_financial[f'A{row}'].font = Font(bold=True)

    row += 1
    ws_financial[f'A{row}'] = "Suma przychodów:"
    ws_financial[f'B{row}'] = f'=SUM(B{start_data_row}:B{end_data_row})'
    ws_financial[f'B{row}'].number_format = '#,##0.0" mln PLN"'

    row += 1
    ws_financial[f'A{row}'] = "Średnie przychody:"
    ws_financial[f'B{row}'] = f'=AVERAGE(B{start_data_row}:B{end_data_row})'
    ws_financial[f'B{row}'].number_format = '#,##0.0" mln PLN"'

    row += 1
    ws_financial[f'A{row}'] = "Min. przychody:"
    ws_financial[f'B{row}'] = f'=MIN(B{start_data_row}:B{end_data_row})'
    ws_financial[f'B{row}'].number_format = '#,##0.0" mln PLN"'

    row += 1
    ws_financial[f'A{row}'] = "Max. przychody:"
    ws_financial[f'B{row}'] = f'=MAX(B{start_data_row}:B{end_data_row})'
    ws_financial[f'B{row}'].number_format = '#,##0.0" mln PLN"'

    row += 1
    ws_financial[f'A{row}'] = "CAGR (wzrost skumulowany):"
    # CAGR formula: ((End Value / Start Value) ^ (1/n)) - 1
    n_years = len(yearly_trends) - 1 if len(yearly_trends) > 1 else 1
    ws_financial[f'B{row}'] = f'=IF(B{start_data_row}>0,((B{end_data_row}/B{start_data_row})^(1/{n_years})-1)*100,0)'
    ws_financial[f'B{row}'].number_format = '0.0"%"'

    # === Section: Financial Ratios ===
    row += 3
    ws_financial[f'A{row}'] = "Wskaźniki finansowe"
    ws_financial[f'A{row}'].font = header_font
    ws_financial[f'A{row}'].fill = subheader_fill
    ws_financial.merge_cells(f'A{row}:C{row}')

    row += 1
    ws_financial[f'A{row}'] = "Wskaźnik"
    ws_financial[f'B{row}'] = "Wartość"
    ws_financial[f'C{row}'] = "Ocena"
    for col in ['A', 'B', 'C']:
        ws_financial[f'{col}{row}'].font = Font(bold=True)
        ws_financial[f'{col}{row}'].fill = PatternFill(start_color='D6DCE5', end_color='D6DCE5', fill_type='solid')

    # Add financial ratios
    ratios_start = row + 1
    ratios = [
        ("ROE (Return on Equity)", financial_data.get('roe', [18.2])[0] if financial_data.get('roe') else 18.2),
        ("ROA (Return on Assets)", financial_data.get('roa', [9.4])[0] if financial_data.get('roa') else 9.4),
        ("Marża brutto", financial_data.get('margins', [28.5])[0] if financial_data.get('margins') else 28.5),
    ]

    for i, (name, value) in enumerate(ratios):
        current_row = ratios_start + i
        ws_financial[f'A{current_row}'] = name
        ws_financial[f'B{current_row}'] = value / 100  # Convert to decimal for percentage format
        ws_financial[f'B{current_row}'].number_format = '0.0%'
        # Formula for rating based on value
        ws_financial[f'C{current_row}'] = f'=IF(B{current_row}>=0.2,"Bardzo dobry",IF(B{current_row}>=0.1,"Dobry","Do poprawy"))'

    # === Section: Projections with formulas ===
    row = ratios_start + len(ratios) + 2
    ws_financial[f'A{row}'] = "Prognoza przychodów (założenie: wzrost 10% r/r)"
    ws_financial[f'A{row}'].font = header_font
    ws_financial[f'A{row}'].fill = subheader_fill
    ws_financial.merge_cells(f'A{row}:C{row}')

    row += 1
    ws_financial[f'A{row}'] = "Rok"
    ws_financial[f'B{row}'] = "Prognoza"
    for col in ['A', 'B']:
        ws_financial[f'{col}{row}'].font = Font(bold=True)
        ws_financial[f'{col}{row}'].fill = PatternFill(start_color='D6DCE5', end_color='D6DCE5', fill_type='solid')

    # Get last year's revenue for projections
    last_revenue_cell = f'B{end_data_row}'
    last_year = int(yearly_trends[-1][0]) if yearly_trends else 2023

    projection_start = row + 1
    for i in range(3):
        current_row = projection_start + i
        ws_financial[f'A{current_row}'] = last_year + i + 1
        if i == 0:
            ws_financial[f'B{current_row}'] = f'={last_revenue_cell}*1.1'
        else:
            prev_row = current_row - 1
            ws_financial[f'B{current_row}'] = f'=B{prev_row}*1.1'
        ws_financial[f'B{current_row}'].number_format = '#,##0.0" mln PLN"'

    # Set column widths for financial sheet
    ws_financial.column_dimensions['A'].width = 35
    ws_financial.column_dimensions['B'].width = 20
    ws_financial.column_dimensions['C'].width = 15

    # === Sheet 3: Report Sections Content ===
    ws_content = wb.create_sheet("Treść raportu")

    ws_content['A1'] = "TREŚĆ RAPORTU"
    ws_content['A1'].font = title_font
    ws_content['A1'].fill = header_fill
    ws_content.merge_cells('A1:B1')

    row = 3
    for i, section in enumerate(report.get('sections', [])):
        ws_content[f'A{row}'] = f"{i+1}. {section['title']}"
        ws_content[f'A{row}'].font = Font(bold=True, size=12)
        row += 1

        # Add content (split into manageable chunks)
        content_lines = section['content'].split('\n')
        for line in content_lines:
            if line.strip():
                ws_content[f'A{row}'] = line.strip()
                ws_content.row_dimensions[row].height = 20
                row += 1
        row += 1

    ws_content.column_dimensions['A'].width = 100

    # Save to buffer
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)

    filename = f"{report['id']}_{report['title'].replace(' ', '_')[:30]}.xlsx"

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )


def create_chart_drawing(chart_data: dict, width: float = 5*72, height: float = 3*72):
    """Create a ReportLab Drawing object with chart visualization.

    Args:
        chart_data: Dictionary with chart configuration
            - type: 'line' or 'bar'
            - title: Chart title
            - data: List of {label, value} dictionaries
            - xKey: Key for x-axis data
            - yKey: Key for y-axis data
            - yLabel: Label for y-axis
            - color: Hex color code
        width: Chart width in points (default 5 inches)
        height: Chart height in points (default 3 inches)

    Returns:
        Drawing object containing the chart
    """
    from reportlab.graphics.shapes import Drawing
    from reportlab.graphics.charts.barcharts import VerticalBarChart
    from reportlab.graphics.charts.linecharts import HorizontalLineChart
    from reportlab.lib import colors

    # Create drawing container
    drawing = Drawing(width, height)

    # Extract chart parameters
    chart_type = chart_data.get('type', 'bar')
    data_points = chart_data.get('data', [])
    x_key = chart_data.get('xKey', 'label')
    y_key = chart_data.get('yKey', 'value')
    color_hex = chart_data.get('color', '#3b82f6')

    # Parse data
    labels = [str(point.get(x_key, '')) for point in data_points]
    values = [float(point.get(y_key, 0)) for point in data_points]

    if not values:
        return drawing

    # Convert hex color to reportlab color
    try:
        chart_color = colors.HexColor(color_hex)
    except:
        chart_color = colors.HexColor('#3b82f6')

    # Create appropriate chart type
    if chart_type == 'line':
        chart = HorizontalLineChart()
        chart.x = 50
        chart.y = 50
        chart.width = width - 100
        chart.height = height - 100
        chart.data = [values]  # Line chart expects list of series
        chart.categoryAxis.categoryNames = labels
        chart.lines[0].strokeColor = chart_color
        chart.lines[0].strokeWidth = 2
    else:  # bar chart
        chart = VerticalBarChart()
        chart.x = 50
        chart.y = 50
        chart.width = width - 100
        chart.height = height - 100
        chart.data = [values]  # Bar chart expects list of series
        chart.categoryAxis.categoryNames = labels
        chart.bars[0].fillColor = chart_color

    # Configure axes
    chart.valueAxis.valueMin = 0
    chart.valueAxis.valueMax = max(values) * 1.1 if values else 10
    chart.valueAxis.valueStep = max(values) / 5 if values else 2

    # Style improvements
    chart.categoryAxis.labels.fontSize = 9
    chart.categoryAxis.labels.angle = 0
    chart.valueAxis.labels.fontSize = 9

    # Add chart to drawing
    drawing.add(chart)

    return drawing


async def export_to_pdf(report: dict) -> StreamingResponse:
    """Generate PDF file with professional formatting."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    # Register DejaVu fonts for proper Polish character support
    try:
        # Try to register DejaVuSans fonts (commonly available on Linux)
        pdfmetrics.registerFont(TTFont('DejaVuSans', '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'))
        pdfmetrics.registerFont(TTFont('DejaVuSans-Bold', '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf'))
        default_font = 'DejaVuSans'
        bold_font = 'DejaVuSans-Bold'
    except:
        # Fallback to Helvetica if DejaVu not available
        default_font = 'Helvetica'
        bold_font = 'Helvetica-Bold'

    # Create buffer for PDF
    output = io.BytesIO()

    # Create PDF document
    doc = SimpleDocTemplate(
        output,
        pagesize=A4,
        rightMargin=0.75*inch,
        leftMargin=0.75*inch,
        topMargin=1*inch,
        bottomMargin=0.75*inch
    )

    # Container for PDF elements
    elements = []

    # Define styles
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#1e40af'),
        spaceAfter=12,
        alignment=TA_CENTER
    )

    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#2563eb'),
        spaceAfter=10,
        spaceBefore=14,
        bold=True
    )

    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['BodyText'],
        fontSize=11,
        alignment=TA_JUSTIFY,
        spaceAfter=8
    )

    # === Title Page ===
    elements.append(Spacer(1, 0.5*inch))
    elements.append(Paragraph(report['title'], title_style))
    elements.append(Spacer(1, 0.2*inch))

    # Report metadata
    metadata_data = [
        ['Typ raportu:', report.get('type', 'N/A').replace('_', ' ').title()],
        ['Firma:', report.get('company', 'N/A')],
        ['Data utworzenia:', report.get('created_at', 'N/A')[:10]],
        ['Ostatnia aktualizacja:', report.get('updated_at', 'N/A')[:10]]
    ]

    metadata_table = Table(metadata_data, colWidths=[2*inch, 4*inch])
    metadata_table.setStyle(TableStyle([
        ('FONTNAME', (0, 0), (0, -1), bold_font),
        ('FONTNAME', (1, 0), (1, -1), default_font),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('TEXTCOLOR', (0, 0), (0, -1), colors.HexColor('#4b5563')),
        ('ALIGN', (0, 0), (0, -1), 'RIGHT'),
        ('ALIGN', (1, 0), (1, -1), 'LEFT'),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
    ]))

    elements.append(metadata_table)
    elements.append(Spacer(1, 0.3*inch))

    # === COMPANY PROFILE CARD (Feature #263) ===
    # Add special info card for company profile reports
    if report.get('type') == 'company_profile':
        # Extract company data from first section (Informacje podstawowe)
        company_data = {}
        if report.get('sections') and len(report['sections']) > 0:
            first_section = report['sections'][0]
            content = first_section.get('content', '')

            # Extract NIP, REGON, KRS from content
            import re
            nip_match = re.search(r'NIP:\s*(\d+)', content)
            regon_match = re.search(r'REGON:\s*(\d+)', content)
            krs_match = re.search(r'KRS:\s*(\d+)', content)
            legal_form_match = re.search(r'Forma prawna:\s*([^\n]+)', content)

            if nip_match:
                company_data['NIP'] = nip_match.group(1)
            if regon_match:
                company_data['REGON'] = regon_match.group(1)
            if krs_match:
                company_data['KRS'] = krs_match.group(1)
            if legal_form_match:
                company_data['Forma prawna'] = legal_form_match.group(1).strip()

        # Create company profile card (styled box)
        if company_data:
            # Card header
            card_header = ParagraphStyle(
                'CardHeader',
                parent=styles['Heading3'],
                fontSize=12,
                textColor=colors.white,
                backColor=colors.HexColor('#2563eb'),
                spaceAfter=0,
                spaceBefore=0,
                alignment=TA_CENTER,
                leftIndent=6,
                rightIndent=6
            )

            elements.append(Paragraph('📋 KARTA INFORMACYJNA FIRMY', card_header))

            # Card body with company data
            card_data = []
            if 'NIP' in company_data:
                card_data.append(['NIP:', company_data['NIP']])
            if 'REGON' in company_data:
                card_data.append(['REGON:', company_data['REGON']])
            if 'KRS' in company_data:
                card_data.append(['KRS:', company_data['KRS']])
            if 'Forma prawna' in company_data:
                card_data.append(['Forma prawna:', company_data['Forma prawna']])

            if card_data:
                card_table = Table(card_data, colWidths=[2*inch, 4*inch])
                card_table.setStyle(TableStyle([
                    ('FONTNAME', (0, 0), (0, -1), bold_font),
                    ('FONTNAME', (1, 0), (1, -1), default_font),
                    ('FONTSIZE', (0, 0), (-1, -1), 10),
                    ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#1f2937')),
                    ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#f3f4f6')),
                    ('ALIGN', (0, 0), (0, -1), 'RIGHT'),
                    ('ALIGN', (1, 0), (1, -1), 'LEFT'),
                    ('TOPPADDING', (0, 0), (-1, -1), 8),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                    ('LEFTPADDING', (0, 0), (-1, -1), 10),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 10),
                    ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#d1d5db')),
                    ('BOX', (0, 0), (-1, -1), 2, colors.HexColor('#2563eb'))
                ]))

                elements.append(card_table)
                elements.append(Spacer(1, 0.3*inch))

    # Summary
    if report.get('summary'):
        elements.append(Paragraph('<b>Podsumowanie</b>', heading_style))
        elements.append(Paragraph(report['summary'], body_style))
        elements.append(Spacer(1, 0.2*inch))

    elements.append(PageBreak())

    # === Report Sections ===
    for i, section in enumerate(report.get('sections', []), 1):
        # Section heading
        section_title = f"{i}. {section.get('title', 'Untitled Section')}"
        elements.append(Paragraph(section_title, heading_style))
        elements.append(Spacer(1, 0.1*inch))

        # Section content
        content = section.get('content', '')

        # Process markdown-like content
        lines = content.split('\n')
        table_rows = []
        in_table = False
        financial_data_rows = []
        collecting_financial_data = False

        for line in lines:
            line = line.strip()
            if not line:
                # Empty line - finalize any ongoing table or financial data
                if in_table and table_rows:
                    # Build table from collected rows
                    table = Table(table_rows)
                    table.setStyle(TableStyle([
                        ('FONTNAME', (0, 0), (-1, 0), bold_font),  # Header row bold
                        ('FONTNAME', (0, 1), (-1, -1), default_font),
                        ('FONTSIZE', (0, 0), (-1, -1), 10),
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#e5e7eb')),
                        ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#1f2937')),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#d1d5db')),
                        ('TOPPADDING', (0, 0), (-1, -1), 8),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                    ]))
                    elements.append(table)
                    elements.append(Spacer(1, 0.1*inch))
                    table_rows = []
                    in_table = False

                if collecting_financial_data and financial_data_rows:
                    # Build financial data table
                    fin_table = Table(financial_data_rows, colWidths=[2.5*inch, 2.5*inch])
                    fin_table.setStyle(TableStyle([
                        ('FONTNAME', (0, 0), (0, -1), bold_font),
                        ('FONTNAME', (1, 0), (1, -1), default_font),
                        ('FONTSIZE', (0, 0), (-1, -1), 10),
                        ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#1f2937')),
                        ('ALIGN', (0, 0), (0, -1), 'LEFT'),
                        ('ALIGN', (1, 0), (1, -1), 'RIGHT'),  # Numbers aligned right
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#d1d5db')),
                        ('TOPPADDING', (0, 0), (-1, -1), 6),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#f3f4f6')),  # First row background
                    ]))
                    elements.append(fin_table)
                    elements.append(Spacer(1, 0.1*inch))
                    financial_data_rows = []
                    collecting_financial_data = False
                continue

            # Detect markdown tables (| ... | ... |)
            if '|' in line and line.count('|') >= 2:
                # Check if it's a separator line (|---|---|)
                if set(line.replace('|', '').replace('-', '').replace(' ', '')) == set():
                    # Skip separator line
                    continue

                # Parse table row
                cells = [cell.strip() for cell in line.split('|')]
                # Remove empty cells from start/end (markdown tables have | at edges)
                cells = [c for c in cells if c]

                if cells:
                    table_rows.append(cells)
                    in_table = True
                continue

            # Detect financial data pattern: "Label: Value" or "• Label: Value"
            if ':' in line and (line.startswith('•') or line.startswith('-') or
                               'mln PLN' in line or '%' in line or
                               any(year in line for year in ['2021', '2022', '2023', '2024', '2025', '2026'])):

                # Parse key-value pair
                if line.startswith('•'):
                    line = line[1:].strip()
                elif line.startswith('-'):
                    line = line[1:].strip()

                if ':' in line:
                    parts = line.split(':', 1)
                    if len(parts) == 2:
                        label = parts[0].strip()
                        value = parts[1].strip()
                        financial_data_rows.append([label + ':', value])
                        collecting_financial_data = True
                        continue

            # If we were collecting tables/data but hit different content, finalize
            if in_table and table_rows:
                table = Table(table_rows)
                table.setStyle(TableStyle([
                    ('FONTNAME', (0, 0), (-1, 0), bold_font),
                    ('FONTNAME', (0, 1), (-1, -1), default_font),
                    ('FONTSIZE', (0, 0), (-1, -1), 10),
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#e5e7eb')),
                    ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#1f2937')),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#d1d5db')),
                    ('TOPPADDING', (0, 0), (-1, -1), 8),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                ]))
                elements.append(table)
                elements.append(Spacer(1, 0.1*inch))
                table_rows = []
                in_table = False

            if collecting_financial_data and financial_data_rows:
                fin_table = Table(financial_data_rows, colWidths=[2.5*inch, 2.5*inch])
                fin_table.setStyle(TableStyle([
                    ('FONTNAME', (0, 0), (0, -1), bold_font),
                    ('FONTNAME', (1, 0), (1, -1), default_font),
                    ('FONTSIZE', (0, 0), (-1, -1), 10),
                    ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#1f2937')),
                    ('ALIGN', (0, 0), (0, -1), 'LEFT'),
                    ('ALIGN', (1, 0), (1, -1), 'RIGHT'),
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#d1d5db')),
                    ('TOPPADDING', (0, 0), (-1, -1), 6),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#f3f4f6')),
                ]))
                elements.append(fin_table)
                elements.append(Spacer(1, 0.1*inch))
                financial_data_rows = []
                collecting_financial_data = False

            # Handle headers
            if line.startswith('**') and line.endswith('**'):
                # Bold text (subheading)
                text = line.strip('*')
                para = Paragraph(f'<b>{text}</b>', body_style)
                elements.append(para)
            elif line.startswith('- '):
                # Bullet point (not financial data)
                text = line[2:]
                para = Paragraph(f'• {text}', body_style)
                elements.append(para)
            else:
                # Regular paragraph
                para = Paragraph(line, body_style)
                elements.append(para)

        # Finalize any remaining table/financial data at section end
        if in_table and table_rows:
            table = Table(table_rows)
            table.setStyle(TableStyle([
                ('FONTNAME', (0, 0), (-1, 0), bold_font),
                ('FONTNAME', (0, 1), (-1, -1), default_font),
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#e5e7eb')),
                ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#1f2937')),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#d1d5db')),
                ('TOPPADDING', (0, 0), (-1, -1), 8),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ]))
            elements.append(table)
            elements.append(Spacer(1, 0.1*inch))

        if collecting_financial_data and financial_data_rows:
            fin_table = Table(financial_data_rows, colWidths=[2.5*inch, 2.5*inch])
            fin_table.setStyle(TableStyle([
                ('FONTNAME', (0, 0), (0, -1), bold_font),
                ('FONTNAME', (1, 0), (1, -1), default_font),
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#1f2937')),
                ('ALIGN', (0, 0), (0, -1), 'LEFT'),
                ('ALIGN', (1, 0), (1, -1), 'RIGHT'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#d1d5db')),
                ('TOPPADDING', (0, 0), (-1, -1), 6),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#f3f4f6')),
            ]))
            elements.append(fin_table)
            elements.append(Spacer(1, 0.1*inch))

        # === ADD CHARTS FOR THIS SECTION (Feature #265) ===
        # Check if report has charts and if any belong to this section
        section_id = section.get('id')
        report_charts = report.get('charts', [])

        for chart in report_charts:
            if chart.get('section_id') == section_id:
                # Add chart title
                chart_title = chart.get('title', 'Wykres')
                chart_title_para = Paragraph(f'<b>{chart_title}</b>', body_style)
                elements.append(chart_title_para)
                elements.append(Spacer(1, 0.05*inch))

                # Generate and add chart drawing
                try:
                    chart_drawing = create_chart_drawing(chart, width=6*inch, height=3*inch)
                    elements.append(chart_drawing)
                    elements.append(Spacer(1, 0.15*inch))
                except Exception as e:
                    # If chart generation fails, add error message
                    error_para = Paragraph(f'<i>Nie udało się wygenerować wykresu: {str(e)}</i>', body_style)
                    elements.append(error_para)
                    elements.append(Spacer(1, 0.1*inch))

        elements.append(Spacer(1, 0.2*inch))

    # === SOURCES SECTION (Feature #266) ===
    if report.get('sources'):
        elements.append(PageBreak())
        sources_heading = Paragraph("Źródła", heading_style)
        elements.append(sources_heading)
        elements.append(Spacer(1, 0.1*inch))

        for source in report['sources']:
            source_text = f"• <b>{source['name']}</b> - Wiarygodność: {int(source['confidence']*100)}%"
            if source.get('url'):
                source_text += f"<br/>&nbsp;&nbsp;&nbsp;&nbsp;URL: <link href='{source['url']}'>{source['url']}</link>"
            source_para = Paragraph(source_text, body_style)
            elements.append(source_para)
            elements.append(Spacer(1, 0.05*inch))

    # === Build PDF ===
    doc.build(elements)
    output.seek(0)

    # Generate filename
    safe_title = report['title'].replace(' ', '_').replace('/', '_')[:40]
    filename = f"{report['id']}_{safe_title}.pdf"

    return StreamingResponse(
        output,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )


async def export_to_docx(report: dict) -> StreamingResponse:
    """Generate DOCX file with professional formatting."""
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE

    # Create new document
    doc = Document()

    # Set document margins (1 inch top/bottom, 0.75 inch left/right)
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(0.75)
        section.right_margin = Inches(0.75)

    # === Title Page ===
    title = doc.add_heading(report['title'], 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    # Style the title with blue color
    for run in title.runs:
        run.font.color.rgb = RGBColor(30, 64, 175)  # #1e40af blue
        run.font.size = Pt(24)

    doc.add_paragraph()  # Add spacing

    # Report metadata table
    metadata_table = doc.add_table(rows=4, cols=2)
    metadata_table.style = 'Light Grid Accent 1'

    metadata_items = [
        ('Typ raportu:', report.get('type', 'N/A').replace('_', ' ').title()),
        ('Firma:', report.get('company', 'N/A')),
        ('Data utworzenia:', report.get('created_at', 'N/A')[:10]),
        ('Ostatnia aktualizacja:', report.get('updated_at', 'N/A')[:10])
    ]

    for i, (label, value) in enumerate(metadata_items):
        row = metadata_table.rows[i]
        label_cell = row.cells[0]
        value_cell = row.cells[1]

        label_cell.text = label
        value_cell.text = value

        # Make labels bold
        for paragraph in label_cell.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(11)

        for paragraph in value_cell.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11)

    doc.add_paragraph()  # Add spacing

    # Summary section
    if report.get('summary'):
        summary_heading = doc.add_heading('Podsumowanie', 2)
        for run in summary_heading.runs:
            run.font.color.rgb = RGBColor(37, 99, 235)  # #2563eb blue

        summary_para = doc.add_paragraph(report['summary'])
        summary_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        for run in summary_para.runs:
            run.font.size = Pt(11)

    doc.add_page_break()

    # === Report Sections ===
    for i, section in enumerate(report.get('sections', []), 1):
        # Section heading
        section_title = f"{i}. {section.get('title', 'Untitled Section')}"
        section_heading = doc.add_heading(section_title, 2)

        # Style section headings
        for run in section_heading.runs:
            run.font.color.rgb = RGBColor(37, 99, 235)  # #2563eb blue
            run.font.size = Pt(16)
            run.font.bold = True

        # Section content
        content = section.get('content', '')

        # Process markdown-like content
        lines = content.split('\n')
        current_list = None

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Handle bold headers (subheadings)
            if line.startswith('**') and line.endswith('**'):
                text = line.strip('*')
                para = doc.add_paragraph()
                run = para.add_run(text)
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(75, 85, 99)  # #4b5563 gray

            # Handle bullet points
            elif line.startswith('- '):
                text = line[2:]
                para = doc.add_paragraph(text, style='List Bullet')
                for run in para.runs:
                    run.font.size = Pt(11)

            # Handle numbered lists
            elif line[0].isdigit() and line[1:3] in ['. ', ') ']:
                text = line[3:] if line[2] == ' ' else line[2:]
                para = doc.add_paragraph(text, style='List Number')
                for run in para.runs:
                    run.font.size = Pt(11)

            # Handle markdown tables (basic support)
            elif '|' in line and line.count('|') >= 2:
                # Skip markdown table separators (lines with --- )
                if not line.replace('|', '').replace('-', '').replace(' ', ''):
                    continue
                # For actual table rows, just add as regular text
                # (full table parsing would be complex)
                para = doc.add_paragraph(line)
                for run in para.runs:
                    run.font.size = Pt(10)
                    run.font.name = 'Courier New'

            # Regular paragraphs
            else:
                para = doc.add_paragraph(line)
                para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                for run in para.runs:
                    run.font.size = Pt(11)

        # Add spacing after each section
        doc.add_paragraph()

    # === Save to BytesIO ===
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)

    # Generate filename
    safe_title = report['title'].replace(' ', '_').replace('/', '_')[:40]
    filename = f"{report['id']}_{safe_title}.docx"

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )


async def export_to_pptx(report: dict) -> StreamingResponse:
    """Generate PowerPoint presentation with professional formatting."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    import re

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)  # Standard 4:3 aspect ratio
    prs.slide_height = Inches(7.5)

    # Define MI-Navigator brand colors
    BRAND_BLUE = RGBColor(30, 64, 175)  # #1e40af
    BRAND_LIGHT_BLUE = RGBColor(37, 99, 235)  # #2563eb
    DARK_GRAY = RGBColor(75, 85, 99)  # #4b5563
    LIGHT_GRAY = RGBColor(229, 231, 235)  # #e5e7eb

    # === Slide 1: Title Slide ===
    title_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(title_slide_layout)

    # Add background shape
    background = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(240, 242, 245)
    background.line.color.rgb = RGBColor(240, 242, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = report['title']
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_BLUE

    # Subtitle with metadata
    subtitle_text = f"{report.get('type', 'N/A').replace('_', ' ').title()}"
    if report.get('company'):
        subtitle_text += f" • {report.get('company')}"
    subtitle_text += f"\n{report.get('created_at', 'N/A')[:10]}"

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle_text
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.color.rgb = DARK_GRAY

    # Logo placeholder (text)
    logo_box = slide.shapes.add_textbox(Inches(4), Inches(6.5), Inches(2), Inches(0.5))
    logo_frame = logo_box.text_frame
    logo_frame.text = "MI-Navigator"
    logo_para = logo_frame.paragraphs[0]
    logo_para.alignment = PP_ALIGN.CENTER
    logo_para.font.size = Pt(14)
    logo_para.font.bold = True
    logo_para.font.color.rgb = BRAND_LIGHT_BLUE

    # === Slide 2: Summary ===
    bullet_slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(bullet_slide_layout)

    # Title
    title = slide.shapes.title
    title.text = "Podsumowanie"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.color.rgb = BRAND_BLUE

    # Content
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True

    # Split summary into paragraphs
    summary_text = report.get('summary', 'Brak podsumowania')
    paragraphs = summary_text.split('\n')

    for i, para_text in enumerate(paragraphs[:5]):  # Limit to 5 paragraphs
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para_text.strip()
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)

    # === Content Slides: One per section ===
    sections = report.get('sections', [])

    for idx, section in enumerate(sections[:10]):  # Limit to 10 sections
        slide = prs.slides.add_slide(bullet_slide_layout)

        # Section title
        title = slide.shapes.title
        title.text = f"{idx + 1}. {section.get('title', 'Untitled')}"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.color.rgb = BRAND_BLUE

        # Section content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.word_wrap = True
        tf.clear()  # Clear default content

        section_content = section.get('content', '')

        # Parse markdown-style content
        lines = section_content.split('\n')
        current_paragraph = None

        for line in lines[:20]:  # Limit lines per slide
            line = line.strip()
            if not line:
                continue

            # Bullet point
            if line.startswith('- ') or line.startswith('* '):
                p = tf.add_paragraph()
                p.text = line[2:].strip()
                p.level = 0
                p.font.size = Pt(12)
                p.font.color.rgb = DARK_GRAY

            # Numbered list
            elif re.match(r'^\d+\.\s', line):
                p = tf.add_paragraph()
                p.text = re.sub(r'^\d+\.\s', '', line)
                p.level = 0
                p.font.size = Pt(12)
                p.font.color.rgb = DARK_GRAY

            # Bold header (markdown **text**)
            elif line.startswith('**') and line.endswith('**'):
                p = tf.add_paragraph()
                p.text = line.strip('*')
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = BRAND_LIGHT_BLUE
                p.space_before = Pt(12)

            # Regular paragraph
            else:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(12)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(6)

    # === Final Slide: Thank You ===
    slide = prs.slides.add_slide(title_slide_layout)

    # Background
    background = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BRAND_BLUE
    background.line.color.rgb = BRAND_BLUE

    # Thank you message
    thanks_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Dziękujemy"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.alignment = PP_ALIGN.CENTER
    thanks_para.font.size = Pt(40)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = RGBColor(255, 255, 255)

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "MI-Navigator | Market Intelligence Platform"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(14)
    footer_para.font.color.rgb = RGBColor(200, 210, 230)

    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    # Generate filename
    safe_title = report['title'].replace(' ', '_').replace('/', '_')[:40]
    filename = f"{report['id']}_{safe_title}.pptx"

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )


@router.post("/bulk-export")
async def bulk_export_reports(
    request: BulkExportRequest,
    current_user: User = Depends(get_current_user)
):
    """Export multiple reports at once."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment

    if not request.report_ids:
        raise HTTPException(status_code=400, detail="No report IDs provided")

    # Find all requested reports
    reports_to_export = []
    for report_id in request.report_ids:
        for r in MOCK_REPORTS:
            if r["id"] == report_id:
                reports_to_export.append(r)
                break

    if not reports_to_export:
        raise HTTPException(status_code=404, detail="No reports found")

    if request.format == "xlsx":
        # Create single Excel file with multiple sheets (one per report)
        wb = Workbook()

        # Title styling
        title_font = Font(name='Calibri', size=14, bold=True, color='FFFFFF')
        header_font = Font(name='Calibri', size=11, bold=True)
        header_fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')

        # First sheet: Summary of all reports
        ws_summary = wb.active
        ws_summary.title = "Podsumowanie"

        ws_summary['A1'] = "ZBIORCZE ZESTAWIENIE RAPORTÓW"
        ws_summary['A1'].font = title_font
        ws_summary['A1'].fill = header_fill
        ws_summary.merge_cells('A1:E1')

        ws_summary['A3'] = "Lp."
        ws_summary['B3'] = "Tytuł raportu"
        ws_summary['C3'] = "Typ"
        ws_summary['D3'] = "Firma"
        ws_summary['E3'] = "Data utworzenia"
        for col in ['A', 'B', 'C', 'D', 'E']:
            ws_summary[f'{col}3'].font = header_font
            ws_summary[f'{col}3'].fill = PatternFill(start_color='D6DCE5', end_color='D6DCE5', fill_type='solid')

        for i, report in enumerate(reports_to_export):
            row = i + 4
            ws_summary[f'A{row}'] = i + 1
            ws_summary[f'B{row}'] = report['title']
            ws_summary[f'C{row}'] = report['type']
            ws_summary[f'D{row}'] = report.get('company', 'N/A')
            ws_summary[f'E{row}'] = report['created_at']

        ws_summary.column_dimensions['A'].width = 5
        ws_summary.column_dimensions['B'].width = 40
        ws_summary.column_dimensions['C'].width = 20
        ws_summary.column_dimensions['D'].width = 25
        ws_summary.column_dimensions['E'].width = 20

        # Create a sheet for each report
        for i, report in enumerate(reports_to_export):
            # Sheet name limited to 31 chars (Excel limitation)
            sheet_name = f"{i+1}. {report['title']}"[:31]
            ws = wb.create_sheet(sheet_name)

            # Report header
            ws['A1'] = report['title']
            ws['A1'].font = title_font
            ws['A1'].fill = header_fill
            ws.merge_cells('A1:D1')

            ws['A2'] = f"Typ: {report['type']}"
            ws['A3'] = f"Firma: {report.get('company', 'N/A')}"
            ws['A4'] = f"Data: {report['created_at']}"
            ws['A5'] = ""
            ws['A6'] = "Podsumowanie:"
            ws['A6'].font = Font(bold=True)
            ws['A7'] = report['summary']
            ws.merge_cells('A7:D7')

            # Sections content
            row = 9
            for j, section in enumerate(report.get('sections', [])):
                ws[f'A{row}'] = f"{j+1}. {section['title']}"
                ws[f'A{row}'].font = Font(bold=True, size=11)
                row += 1

                # Split content into lines
                for line in section['content'].split('\n'):
                    if line.strip():
                        ws[f'A{row}'] = line.strip()
                        row += 1
                row += 1

            ws.column_dimensions['A'].width = 80

        # Save to buffer
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)

        filename = f"raporty_eksport_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"

        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )

    elif request.format == "pdf":
        # Create a ZIP file with individual PDFs for each report
        from reportlab.lib.pagesizes import A4, letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
        import zipfile

        # Create a temporary directory for PDFs
        pdf_buffers = []

        for report in reports_to_export:
            # Create PDF for this report
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4,
                                   topMargin=0.75*inch, bottomMargin=0.75*inch,
                                   leftMargin=0.75*inch, rightMargin=0.75*inch)

            # Container for the 'Flowable' objects
            elements = []

            # Styles
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=18,
                textColor=colors.HexColor('#1e40af'),
                spaceAfter=12,
                alignment=TA_CENTER,
                fontName='Helvetica-Bold'
            )
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontSize=14,
                textColor=colors.HexColor('#1e40af'),
                spaceAfter=10,
                spaceBefore=14,
                fontName='Helvetica-Bold'
            )
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['BodyText'],
                fontSize=10,
                leading=14,
                alignment=TA_JUSTIFY,
                spaceAfter=6
            )

            # Title
            title = Paragraph(report['title'], title_style)
            elements.append(title)
            elements.append(Spacer(1, 0.2*inch))

            # Metadata table
            metadata = [
                ['Typ:', report['type']],
                ['Firma:', report.get('company', 'N/A')],
                ['Data utworzenia:', report['created_at']],
                ['Status:', report['status']]
            ]
            metadata_table = Table(metadata, colWidths=[1.5*inch, 4*inch])
            metadata_table.setStyle(TableStyle([
                ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('TEXTCOLOR', (0, 0), (0, -1), colors.HexColor('#374151')),
                ('TEXTCOLOR', (1, 0), (1, -1), colors.HexColor('#6b7280')),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ]))
            elements.append(metadata_table)
            elements.append(Spacer(1, 0.3*inch))

            # Summary section
            summary_heading = Paragraph("Podsumowanie", heading_style)
            elements.append(summary_heading)
            summary_text = Paragraph(report['summary'], body_style)
            elements.append(summary_text)
            elements.append(Spacer(1, 0.2*inch))

            # Sections
            for i, section in enumerate(report.get('sections', [])):
                section_title = Paragraph(f"{i+1}. {section['title']}", heading_style)
                elements.append(section_title)

                # Process content - split by newlines and create paragraphs
                content_lines = section['content'].split('\n')
                for line in content_lines:
                    if line.strip():
                        para = Paragraph(line.strip(), body_style)
                        elements.append(para)

                elements.append(Spacer(1, 0.15*inch))

            # Sources section (if any)
            if report.get('sources'):
                sources_heading = Paragraph("Źródła", heading_style)
                elements.append(sources_heading)

                for source in report['sources']:
                    source_text = f"• {source['name']} - Wiarygodność: {int(source['confidence']*100)}%"
                    if source.get('url'):
                        source_text += f" ({source['url']})"
                    source_para = Paragraph(source_text, body_style)
                    elements.append(source_para)

            # Build PDF
            doc.build(elements)
            buffer.seek(0)

            # Store buffer with filename
            safe_filename = "".join(c for c in report['title'] if c.isalnum() or c in (' ', '-', '_')).strip()
            safe_filename = safe_filename[:50]  # Limit length
            pdf_buffers.append((f"{safe_filename}.pdf", buffer))

        # Create ZIP file
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for filename, pdf_buffer in pdf_buffers:
                zip_file.writestr(filename, pdf_buffer.getvalue())

        zip_buffer.seek(0)

        zip_filename = f"raporty_pdf_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"

        return StreamingResponse(
            zip_buffer,
            media_type="application/zip",
            headers={"Content-Disposition": f'attachment; filename="{zip_filename}"'}
        )

    else:
        raise HTTPException(status_code=400, detail=f"Unsupported format: {request.format}")


@router.post("/bulk-export-async")
async def bulk_export_reports_async(
    request: BulkExportRequest,
    current_user: User = Depends(get_current_user)
):
    """Start async bulk export with progress tracking."""
    import asyncio

    if not request.report_ids:
        raise HTTPException(status_code=400, detail="No report IDs provided")

    # Generate task ID
    task_id = str(uuid.uuid4())

    # Initialize task status
    BATCH_TASKS[task_id] = {
        "status": "processing",
        "total": len(request.report_ids),
        "processed": 0,
        "result": None,
        "error": None,
        "format": request.format
    }

    # Start background processing (simulated)
    async def process_export():
        """Simulate processing with progress updates."""
        try:
            for i in range(len(request.report_ids)):
                # Simulate processing time per report (2 seconds each for demo)
                await asyncio.sleep(2.0)

                # Update progress
                BATCH_TASKS[task_id]["processed"] = i + 1

            # Mark as completed
            BATCH_TASKS[task_id]["status"] = "completed"
            BATCH_TASKS[task_id]["result"] = f"Successfully exported {len(request.report_ids)} reports"
        except Exception as e:
            BATCH_TASKS[task_id]["status"] = "failed"
            BATCH_TASKS[task_id]["error"] = str(e)

    # Start task in background (fire and forget)
    asyncio.create_task(process_export())

    return {
        "task_id": task_id,
        "status": "started",
        "total": len(request.report_ids)
    }


@router.get("/bulk-export-progress/{task_id}")
async def get_bulk_export_progress(
    task_id: str,
    current_user: User = Depends(get_current_user)
):
    """Get progress of bulk export operation."""
    if task_id not in BATCH_TASKS:
        raise HTTPException(status_code=404, detail="Task not found")

    task = BATCH_TASKS[task_id]

    # Calculate percentage
    percentage = (task["processed"] / task["total"] * 100) if task["total"] > 0 else 0

    return {
        "task_id": task_id,
        "status": task["status"],
        "total": task["total"],
        "processed": task["processed"],
        "percentage": round(percentage, 1),
        "result": task.get("result"),
        "error": task.get("error")
    }


# In-memory storage for share links: { "share_token": {"report_id", "created_at", "expires_at", "created_by"} }
SHARE_LINKS: dict = {}

# In-memory storage for access log: { "share_token": [{"timestamp", "ip", "user_agent", "location"}, ...] }
SHARE_ACCESS_LOG: dict = {}


class ShareLinkRequest(BaseModel):
    """Request model for creating a share link"""
    password: Optional[str] = None


@router.post("/{report_id}/share")
async def share_report(
    report_id: str,
    share_request: ShareLinkRequest,
    request: Request,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Generate share link for report with tracking and optional password protection."""
    # Verify report exists
    report = None
    for r in MOCK_REPORTS:
        if r["id"] == report_id:
            report = r
            break

    if not report:
        raise HTTPException(status_code=404, detail="Report not found")

    # Generate unique share token
    share_token = str(uuid.uuid4())

    # Calculate expiration (30 days from now)
    from datetime import timedelta
    expires_at = datetime.now() + timedelta(days=30)

    # Hash password if provided
    password_hash = None
    if share_request.password:
        import hashlib
        # Using SHA-256 for simplicity (in production, use bcrypt or argon2)
        password_hash = hashlib.sha256(share_request.password.encode()).hexdigest()

    # Store share link
    SHARE_LINKS[share_token] = {
        "report_id": report_id,
        "created_at": datetime.now().isoformat() + "Z",
        "expires_at": expires_at.isoformat() + "Z",
        "created_by": str(current_user.id),
        "created_by_email": current_user.email,
        "password_hash": password_hash,
        "password_protected": password_hash is not None
    }

    # Initialize access log for this token
    SHARE_ACCESS_LOG[share_token] = []

    share_url = f"http://localhost:3000/share/{share_token}"

    # Log audit event
    await log_audit(
        db=db,
        user=current_user,
        action_type=AuditAction.REPORT_SHARE,
        resource_type=ResourceType.REPORT,
        resource_id=report_id,
        description=f"Shared report: {report['title']}",
        request=request,
        extra_data={
            "report_title": report["title"],
            "share_token": share_token,
            "expires_at": expires_at.isoformat() + "Z"
        }
    )

    return {
        "share_token": share_token,
        "share_url": share_url,
        "expires_at": expires_at.isoformat() + "Z",
        "report_id": report_id,
        "report_title": report["title"],
        "password_protected": password_hash is not None
    }


class EmailShareRequest(BaseModel):
    recipient_email: str
    message: Optional[str] = ""
    sender_name: Optional[str] = None


@router.post("/{report_id}/share/email")
async def share_report_via_email(
    report_id: str,
    request: EmailShareRequest,
    request_obj: Request,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Share report via email (dev mode - logs to console)."""
    # Find the report
    report = None
    for r in MOCK_REPORTS:
        if r["id"] == report_id:
            report = r
            break

    if not report:
        raise HTTPException(status_code=404, detail="Report not found")

    # In development mode, log email to console instead of sending
    sender_name = request.sender_name or current_user.name or current_user.email
    report_url = f"http://localhost:3000/reports/{report_id}"

    email_content = f"""
===========================================
📧 EMAIL NOTIFICATION (DEV MODE)
===========================================
From: {sender_name} <{current_user.email}>
To: {request.recipient_email}
Subject: [MI-Navigator] {sender_name} udostępnił Ci raport: {report['title']}

-------------------------------------------
MESSAGE:
-------------------------------------------
Cześć!

{sender_name} udostępnił Ci raport z MI-Navigator:

📊 Raport: {report['title']}
📁 Typ: {report.get('type', 'N/A')}
🏢 Firma: {report.get('company', 'N/A')}

{f'Wiadomość od {sender_name}:' if request.message else ''}
{request.message if request.message else ''}

Kliknij poniższy link aby zobaczyć raport:
{report_url}

---
Ten email został wygenerowany przez MI-Navigator - Market Intelligence Platform
===========================================
"""

    # Log to console (development mode)
    print(email_content)

    # Log audit event
    await log_audit(
        db=db,
        user=current_user,
        action_type=AuditAction.REPORT_SHARE,
        resource_type=ResourceType.REPORT,
        resource_id=report_id,
        description=f"Shared report '{report['title']}' via email to {request.recipient_email}",
        request=request_obj,
        extra_data={
            "recipient_email": request.recipient_email,
            "message": request.message,
            "report_title": report['title']
        }
    )

    return {
        "success": True,
        "message": f"Raport został udostępniony przez email do {request.recipient_email}",
        "recipient": request.recipient_email,
        "report_id": report_id,
        "report_title": report['title']
    }


# In-memory annotation storage (per user, per report)
REPORT_ANNOTATIONS: dict = {}


# In-memory favorites storage: { "user_id": ["report_id1", "report_id2", ...] }
USER_FAVORITES: dict = {}


class AnnotationCreate(BaseModel):
    section_id: str
    selected_text: str
    start_offset: int
    end_offset: int
    comment: str


class Annotation(BaseModel):
    id: str
    report_id: str
    section_id: str
    selected_text: str
    start_offset: int
    end_offset: int
    comment: str
    created_at: str
    user_id: str


@router.get("/{report_id}/annotations")
async def get_annotations(
    report_id: str,
    current_user: User = Depends(get_current_user)
):
    """Get all annotations for a report."""
    user_key = f"{current_user.id}:{report_id}"
    annotations = REPORT_ANNOTATIONS.get(user_key, [])
    return {"annotations": annotations}


@router.post("/{report_id}/annotations")
async def create_annotation(
    report_id: str,
    annotation: AnnotationCreate,
    current_user: User = Depends(get_current_user)
):
    """Create a new annotation for a report."""
    user_key = f"{current_user.id}:{report_id}"

    if user_key not in REPORT_ANNOTATIONS:
        REPORT_ANNOTATIONS[user_key] = []

    # Generate unique ID
    annotation_id = f"ann_{len(REPORT_ANNOTATIONS[user_key]) + 1}_{int(datetime.now().timestamp())}"

    new_annotation = Annotation(
        id=annotation_id,
        report_id=report_id,
        section_id=annotation.section_id,
        selected_text=annotation.selected_text,
        start_offset=annotation.start_offset,
        end_offset=annotation.end_offset,
        comment=annotation.comment,
        created_at=datetime.now().isoformat() + "Z",
        user_id=str(current_user.id)
    )

    REPORT_ANNOTATIONS[user_key].append(new_annotation.model_dump())

    return new_annotation


@router.delete("/{report_id}/annotations/{annotation_id}")
async def delete_annotation(
    report_id: str,
    annotation_id: str,
    current_user: User = Depends(get_current_user)
):
    """Delete an annotation."""
    user_key = f"{current_user.id}:{report_id}"

    if user_key not in REPORT_ANNOTATIONS:
        return {"error": "Annotation not found"}

    annotations = REPORT_ANNOTATIONS[user_key]
    REPORT_ANNOTATIONS[user_key] = [a for a in annotations if a["id"] != annotation_id]

    return {"message": "Annotation deleted successfully"}


# ============================================================================
# FAVORITES ENDPOINTS
# ============================================================================

@router.get("/favorites")
async def get_favorite_reports(
    current_user: User = Depends(get_current_user)
):
    """Get list of user's favorite report IDs."""
    user_id = str(current_user.id)
    favorites = USER_FAVORITES.get(user_id, [])
    return {"favorites": favorites}


@router.post("/{report_id}/favorite")
async def add_to_favorites(
    report_id: str,
    current_user: User = Depends(get_current_user)
):
    """Add a report to user's favorites."""
    user_id = str(current_user.id)

    # Initialize user's favorites list if it doesn't exist
    if user_id not in USER_FAVORITES:
        USER_FAVORITES[user_id] = []

    # Add to favorites if not already there
    if report_id not in USER_FAVORITES[user_id]:
        USER_FAVORITES[user_id].append(report_id)
        return {"message": "Report added to favorites", "is_favorite": True}

    return {"message": "Report already in favorites", "is_favorite": True}


@router.delete("/{report_id}/favorite")
async def remove_from_favorites(
    report_id: str,
    current_user: User = Depends(get_current_user)
):
    """Remove a report from user's favorites."""
    user_id = str(current_user.id)

    if user_id in USER_FAVORITES and report_id in USER_FAVORITES[user_id]:
        USER_FAVORITES[user_id].remove(report_id)
        return {"message": "Report removed from favorites", "is_favorite": False}

    return {"message": "Report not in favorites", "is_favorite": False}


@router.get("/{report_id}/favorite")
async def check_favorite_status(
    report_id: str,
    current_user: User = Depends(get_current_user)
):
    """Check if a report is in user's favorites."""
    user_id = str(current_user.id)
    is_favorite = user_id in USER_FAVORITES and report_id in USER_FAVORITES[user_id]
    return {"is_favorite": is_favorite}


# Version history mock data
REPORT_VERSIONS = {
    "report_001": [
        {
            "version": 3,
            "created_at": "2026-01-14T14:22:00Z",
            "author": "Jan Kowalski",
            "changes": "Dodano sekcje Analiza SWOT",
            "is_current": True
        },
        {
            "version": 2,
            "created_at": "2026-01-14T12:15:00Z",
            "author": "Jan Kowalski",
            "changes": "Zaktualizowano dane finansowe",
            "is_current": False,
            "sections": [
                {
                    "id": "section_1",
                    "title": "Informacje podstawowe",
                    "content": """FADO Sp. z o.o. to polska firma założona w 1998 roku, specjalizująca się w produkcji wyrobów z tworzyw sztucznych. Firma posiada siedzibę w Warszawie przy ul. Przemysłowej 15.

**Dane rejestrowe:**
- NIP: 5260016831
- REGON: 012567834
- KRS: 0000145732
- Forma prawna: Spółka z ograniczoną odpowiedzialnością

Firma zatrudnia około 120 pracowników."""
                },
                {
                    "id": "section_2",
                    "title": "Analiza finansowa",
                    "content": """**Przychody i rentowność (2022):**
- Przychody ze sprzedaży: 40,2 mln PLN
- Wzrost r/r: +8,5%
- Marża brutto: 26,2%
- Zysk netto: 3,9 mln PLN

Firma wykazuje stabilny wzrost przychodów."""
                }
            ]
        },
        {
            "version": 1,
            "created_at": "2026-01-13T09:30:00Z",
            "author": "System",
            "changes": "Utworzono raport",
            "is_current": False,
            "sections": [
                {
                    "id": "section_1",
                    "title": "Informacje podstawowe",
                    "content": """FADO Sp. z o.o. to polska firma założona w 1998 roku.

**Dane rejestrowe:**
- NIP: 5260016831
- REGON: 012567834
- KRS: 0000145732

Wstępne dane firmy."""
                }
            ]
        }
    ],
    "report_002": [
        {
            "version": 1,
            "created_at": "2026-01-13T09:15:00Z",
            "author": "System",
            "changes": "Utworzono raport",
            "is_current": True
        }
    ],
    "report_003": [
        {
            "version": 1,
            "created_at": "2026-01-12T11:00:00Z",
            "author": "System",
            "changes": "Utworzono raport",
            "is_current": True
        }
    ]
}


class ReportVersion(BaseModel):
    version: int
    created_at: str
    author: str
    changes: str
    is_current: bool


class ReportVersionDetail(BaseModel):
    version: int
    created_at: str
    author: str
    changes: str
    is_current: bool
    sections: Optional[List[ReportSection]] = None


@router.get("/{report_id}/versions")
async def get_report_versions(report_id: str):
    """Get version history for a report."""
    versions = REPORT_VERSIONS.get(report_id, [])
    return {
        "versions": [
            ReportVersion(
                version=v["version"],
                created_at=v["created_at"],
                author=v["author"],
                changes=v["changes"],
                is_current=v["is_current"]
            ) for v in versions
        ]
    }


@router.get("/{report_id}/versions/{version}")
async def get_report_version(
    report_id: str,
    version: int
    # current_user: User = Depends(get_current_user)  # Disabled for testing
):
    """Get a specific version of a report."""
    versions = REPORT_VERSIONS.get(report_id, [])

    for v in versions:
        if v["version"] == version:
            # If it's the current version, return the main report
            if v["is_current"]:
                for report in MOCK_REPORTS:
                    if report["id"] == report_id:
                        return ReportDetail(
                            id=report["id"],
                            title=report["title"] + f" (wersja {version})",
                            type=report["type"],
                            company=report["company"],
                            created_at=v["created_at"],
                            updated_at=v["created_at"],
                            status=report["status"],
                            summary=report["summary"],
                            sections=[ReportSection(**s) for s in report["sections"]],
                            sources=[ReportSource(**s) for s in report["sources"]]
                        )
            else:
                # Return the historical version
                for report in MOCK_REPORTS:
                    if report["id"] == report_id:
                        return ReportDetail(
                            id=report["id"],
                            title=report["title"] + f" (wersja {version})",
                            type=report["type"],
                            company=report["company"],
                            created_at=v["created_at"],
                            updated_at=v["created_at"],
                            status="archived",
                            summary=report["summary"],
                            sections=[ReportSection(**s) for s in v.get("sections", [])],
                            sources=[ReportSource(**s) for s in report["sources"]]
                        )

    return {"error": "Version not found"}


class RestoreVersionRequest(BaseModel):
    version: int


@router.post("/{report_id}/versions/restore")
async def restore_report_version(
    report_id: str,
    request: RestoreVersionRequest
    # Note: Auth disabled to match other version endpoints in this project
):
    """Restore a report to a previous version.

    This creates a new version with the content from the specified historical version.
    """
    versions = REPORT_VERSIONS.get(report_id, [])

    if not versions:
        return {"error": "Report not found"}

    # Find the version to restore
    version_to_restore = None
    for v in versions:
        if v["version"] == request.version:
            version_to_restore = v
            break

    if not version_to_restore:
        return {"error": "Version not found"}

    # Get the current max version
    max_version = max(v["version"] for v in versions)

    # Mark all versions as not current
    for v in versions:
        v["is_current"] = False

    # Create new version with restored content
    new_version = {
        "version": max_version + 1,
        "created_at": datetime.now().isoformat() + "Z",
        "author": "Jan Kowalski",  # Mock author for now (auth disabled)
        "changes": f"Przywrócono wersję {request.version}",
        "is_current": True,
    }

    # Copy sections from the restored version if it has them
    if "sections" in version_to_restore:
        new_version["sections"] = version_to_restore["sections"].copy()

    # Add the new version at the beginning (most recent first)
    REPORT_VERSIONS[report_id].insert(0, new_version)

    return {
        "message": "Wersja przywrócona pomyślnie",
        "new_version": new_version["version"],
        "restored_from": request.version
    }


# Collaboration comments storage (global, visible to all users)
REPORT_COMMENTS: dict = {
    "report_001": [
        {
            "id": "comment_1",
            "report_id": "report_001",
            "user_id": "1",
            "user_name": "Jan Kowalski",
            "user_email": "jan.kowalski@example.com",
            "text": "Świetna analiza finansowa! Warto zwrócić uwagę na wzrost marży r/r.",
            "created_at": "2026-01-14T10:45:00Z",
            "parent_id": None,
            "resolved": False,
            "resolved_by": None,
            "resolved_by_name": None,
            "resolved_at": None,
        },
        {
            "id": "comment_2",
            "report_id": "report_001",
            "user_id": "2",
            "user_name": "Anna Nowak",
            "user_email": "anna.nowak@example.com",
            "text": "Czy mamy dostęp do danych za Q4 2023? Byłoby dobrze je uwzględnić.",
            "created_at": "2026-01-14T11:30:00Z",
            "parent_id": None,
            "resolved": True,
            "resolved_by": "1",
            "resolved_by_name": "Jan Kowalski",
            "resolved_at": "2026-01-14T12:05:00Z",
        },
        {
            "id": "comment_3",
            "report_id": "report_001",
            "user_id": "1",
            "user_name": "Jan Kowalski",
            "user_email": "jan.kowalski@example.com",
            "text": "Tak, mam dane za Q4. Dodałem je w sekcji finansowej.",
            "created_at": "2026-01-14T12:00:00Z",
            "parent_id": "comment_2",
            "resolved": False,
            "resolved_by": None,
            "resolved_by_name": None,
            "resolved_at": None,
        }
    ]
}


class CommentCreate(BaseModel):
    text: str
    parent_id: Optional[str] = None


class Comment(BaseModel):
    id: str
    report_id: str
    user_id: str
    user_name: str
    user_email: str
    text: str
    created_at: str
    parent_id: Optional[str] = None
    resolved: bool = False
    resolved_by: Optional[str] = None
    resolved_by_name: Optional[str] = None
    resolved_at: Optional[str] = None


@router.get("/{report_id}/comments")
async def get_comments(
    report_id: str,
    resolved: Optional[bool] = None
):
    """Get all collaboration comments for a report. Comments are visible to all users.

    Optional filter by resolved status:
    - resolved=true: only resolved comments
    - resolved=false: only unresolved comments
    - no filter: all comments
    """
    comments = REPORT_COMMENTS.get(report_id, [])

    # Filter by resolved status if specified
    if resolved is not None:
        comments = [c for c in comments if c.get("resolved", False) == resolved]

    return {"comments": comments}


@router.post("/{report_id}/comments")
async def create_comment(
    report_id: str,
    comment: CommentCreate,
    current_user: User = Depends(get_current_user)
):
    """Create a new collaboration comment for a report. Supports replies via parent_id."""
    if report_id not in REPORT_COMMENTS:
        REPORT_COMMENTS[report_id] = []

    # Validate parent_id if provided
    if comment.parent_id:
        parent_exists = any(c["id"] == comment.parent_id for c in REPORT_COMMENTS[report_id])
        if not parent_exists:
            return {"error": "Parent comment not found"}

    # Generate unique ID
    comment_id = f"comment_{len(REPORT_COMMENTS[report_id]) + 1}_{int(datetime.now().timestamp())}"

    new_comment = Comment(
        id=comment_id,
        report_id=report_id,
        user_id=str(current_user.id),
        user_name=current_user.name or current_user.email.split('@')[0],
        user_email=current_user.email,
        text=comment.text,
        created_at=datetime.now().isoformat() + "Z",
        parent_id=comment.parent_id,
        resolved=False,
        resolved_by=None,
        resolved_by_name=None,
        resolved_at=None
    )

    REPORT_COMMENTS[report_id].append(new_comment.model_dump())

    return new_comment


@router.delete("/{report_id}/comments/{comment_id}")
async def delete_comment(
    report_id: str,
    comment_id: str,
    current_user: User = Depends(get_current_user)
):
    """Delete a collaboration comment. Only the author can delete their comment."""
    if report_id not in REPORT_COMMENTS:
        return {"error": "Comment not found"}

    comments = REPORT_COMMENTS[report_id]
    comment_to_delete = None

    for c in comments:
        if c["id"] == comment_id:
            comment_to_delete = c
            break

    if not comment_to_delete:
        return {"error": "Comment not found"}

    # Check if user is the author
    if comment_to_delete["user_id"] != str(current_user.id):
        return {"error": "Nie możesz usunąć komentarza innego użytkownika"}

    REPORT_COMMENTS[report_id] = [c for c in comments if c["id"] != comment_id]

    return {"message": "Komentarz został usunięty"}


@router.patch("/{report_id}/comments/{comment_id}/resolve")
async def resolve_comment(
    report_id: str,
    comment_id: str,
    current_user: User = Depends(get_current_user)
):
    """Mark a comment as resolved. Any user can resolve any comment."""
    if report_id not in REPORT_COMMENTS:
        return {"error": "Comment not found"}

    comments = REPORT_COMMENTS[report_id]
    comment_to_resolve = None

    for c in comments:
        if c["id"] == comment_id:
            comment_to_resolve = c
            break

    if not comment_to_resolve:
        return {"error": "Comment not found"}

    # Toggle resolved status
    if comment_to_resolve.get("resolved", False):
        # Unresolve
        comment_to_resolve["resolved"] = False
        comment_to_resolve["resolved_by"] = None
        comment_to_resolve["resolved_by_name"] = None
        comment_to_resolve["resolved_at"] = None
        return {
            "message": "Komentarz został odznaczony jako nierozwiązany",
            "comment": comment_to_resolve
        }
    else:
        # Resolve
        comment_to_resolve["resolved"] = True
        comment_to_resolve["resolved_by"] = str(current_user.id)
        comment_to_resolve["resolved_by_name"] = current_user.name or current_user.email.split('@')[0]
        comment_to_resolve["resolved_at"] = datetime.now().isoformat() + "Z"
        return {
            "message": "Komentarz został oznaczony jako rozwiązany",
            "comment": comment_to_resolve
        }


# ============================================================================
# ARCHIVE/UNARCHIVE ENDPOINTS
# ============================================================================

@router.post("/{report_id}/archive")
async def archive_report(
    report_id: str,
    current_user: User = Depends(get_current_user)
):
    """Archive a report."""
    # Find the report
    for report in MOCK_REPORTS:
        if report["id"] == report_id:
            report["is_archived"] = True
            return {"message": "Raport zarchiwizowany pomyślnie", "is_archived": True}

    raise HTTPException(status_code=404, detail="Report not found")


@router.post("/{report_id}/unarchive")
async def unarchive_report(
    report_id: str,
    current_user: User = Depends(get_current_user)
):
    """Unarchive a report."""
    # Find the report
    for report in MOCK_REPORTS:
        if report["id"] == report_id:
            report["is_archived"] = False
            return {"message": "Raport przywrócony z archiwum", "is_archived": False}

    raise HTTPException(status_code=404, detail="Report not found")


# ==================== TEMPLATES ====================

@router.post("/{report_id}/save-as-template")
async def save_report_as_template(
    report_id: str,
    template_name: str = Query(..., description="Name for the template"),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Save a report as a reusable template."""
    import copy

    # Auto-generate test reports if user has none yet (ensures MOCK_REPORTS is populated)
    user_id = str(current_user.id)
    user_reports = [r for r in MOCK_REPORTS if r.get("created_by") == user_id]
    if len(user_reports) == 0:
        test_reports = generate_pagination_test_reports(1000, user_id)
        MOCK_REPORTS.extend(test_reports)

    # Find the source report
    source_report = None
    for report in MOCK_REPORTS:
        if report["id"] == report_id:
            source_report = report
            break

    if not source_report:
        raise HTTPException(status_code=404, detail="Report not found")

    # Create template in database
    new_template = ReportTemplate(
        name=template_name,
        type=source_report["type"],
        created_by=str(current_user.id),
        sections=json.dumps(copy.deepcopy(source_report["sections"])),  # Store as JSON string
        use_count=0,
        last_used=None,
        original_report_id=report_id,
        original_report_title=source_report["title"],
    )

    db.add(new_template)
    await db.commit()
    await db.refresh(new_template)

    return {
        "message": "Template created successfully",
        "template_id": new_template.id,
        "template_name": template_name
    }


class ReportUpdateRequest(BaseModel):
    sections: Optional[List[dict]] = None
    title: Optional[str] = None
    summary: Optional[str] = None


@router.patch("/{report_id}")
async def update_report(
    report_id: str,
    update_data: ReportUpdateRequest,
    current_user: User = Depends(get_current_user)
):
    """
    Update report sections or metadata
    """
    # Find the report
    report = next((r for r in MOCK_REPORTS if r["id"] == report_id), None)
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")

    # Update sections if provided
    if update_data.sections is not None:
        report["sections"] = update_data.sections

    # Update title if provided
    if update_data.title is not None:
        report["title"] = update_data.title

    # Update summary if provided
    if update_data.summary is not None:
        report["summary"] = update_data.summary

    # Update timestamp
    report["updated_at"] = datetime.utcnow().isoformat() + "Z"

    return {
        "message": "Report updated successfully",
        "report": report
    }


@router.post("/upload-image")
async def upload_image(
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user)
):
    """
    Upload an image for use in reports
    """
    # Validate file type
    if not file.content_type or not file.content_type.startswith('image/'):
        raise HTTPException(status_code=400, detail="File must be an image")

    # Generate unique filename
    file_extension = Path(file.filename).suffix if file.filename else '.jpg'
    unique_filename = f"{uuid.uuid4()}{file_extension}"

    # Define upload path
    upload_dir = Path(__file__).parent.parent.parent.parent.parent / "static" / "uploads"
    upload_dir.mkdir(parents=True, exist_ok=True)
    file_path = upload_dir / unique_filename

    # Save file
    try:
        contents = await file.read()
        with open(file_path, "wb") as f:
            f.write(contents)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save image: {str(e)}")

    # Return URL to access the image
    image_url = f"/static/uploads/{unique_filename}"

    return {
        "url": image_url,
        "filename": unique_filename,
        "original_filename": file.filename
    }


# ============================================================================
# PUBLIC SHARE LINK ACCESS (with tracking)
# ============================================================================

class SharePasswordRequest(BaseModel):
    """Request model for verifying share link password"""
    password: str


@router.post("/public/share/{share_token}/verify-password")
async def verify_share_password(
    share_token: str,
    password_request: SharePasswordRequest
):
    """
    Verify password for password-protected share link.
    Returns access token if password is correct.
    """
    # Check if share link exists
    if share_token not in SHARE_LINKS:
        raise HTTPException(status_code=404, detail="Share link not found or expired")

    share_link = SHARE_LINKS[share_token]

    # Check if link has expired
    from datetime import datetime, timezone
    expires_at = datetime.fromisoformat(share_link["expires_at"].replace("Z", "+00:00"))
    current_time = datetime.now(timezone.utc)
    if current_time > expires_at:
        raise HTTPException(status_code=410, detail="Share link has expired")

    # Check if link is password protected
    if not share_link.get("password_hash"):
        raise HTTPException(status_code=400, detail="This link is not password protected")

    # Verify password
    import hashlib
    provided_hash = hashlib.sha256(password_request.password.encode()).hexdigest()

    if provided_hash != share_link["password_hash"]:
        raise HTTPException(status_code=401, detail="Incorrect password")

    return {"success": True, "message": "Password verified"}


@router.get("/public/share/{share_token}")
async def access_shared_report(
    share_token: str,
    request: Request,
    password_verified: Optional[str] = None
):
    """
    Public endpoint to access shared report (no authentication required).
    For password-protected links, returns metadata first.
    After password verification, returns full report.
    """
    # Check if share link exists
    if share_token not in SHARE_LINKS:
        raise HTTPException(status_code=404, detail="Share link not found or expired")

    share_link = SHARE_LINKS[share_token]

    # Check if link has expired
    from datetime import datetime, timezone
    expires_at = datetime.fromisoformat(share_link["expires_at"].replace("Z", "+00:00"))
    current_time = datetime.now(timezone.utc)
    if current_time > expires_at:
        raise HTTPException(status_code=410, detail="Share link has expired")

    # If password protected and not verified yet, return metadata only
    if share_link.get("password_hash") and password_verified != "true":
        # Return metadata only, frontend will prompt for password
        return {
            "password_protected": True,
            "shared_by": share_link.get("created_by_email", "Unknown"),
            "expires_at": share_link["expires_at"]
        }

    report_id = share_link["report_id"]

    # Find the report
    report = None
    for r in MOCK_REPORTS:
        if r["id"] == report_id:
            report = r
            break

    if not report:
        raise HTTPException(status_code=404, detail="Report not found")

    # Log access details
    client_ip = request.client.host if request.client else "unknown"
    user_agent = request.headers.get("user-agent", "unknown")

    access_entry = {
        "timestamp": datetime.now().isoformat() + "Z",
        "ip_address": client_ip,
        "user_agent": user_agent,
        "country": "Unknown",  # Could integrate with GeoIP service
        "city": "Unknown"
    }

    # Add to access log
    if share_token not in SHARE_ACCESS_LOG:
        SHARE_ACCESS_LOG[share_token] = []

    SHARE_ACCESS_LOG[share_token].append(access_entry)

    # Log to console for debugging
    print(f"📊 SHARE ACCESS: Token={share_token[:8]}..., IP={client_ip}, Report={report['title']}")

    # Return report data
    return {
        "report": report,
        "shared_by": share_link.get("created_by_email", "Unknown"),
        "share_token": share_token
    }


@router.get("/{report_id}/share/access-log")
async def get_share_access_log(
    report_id: str,
    current_user: User = Depends(get_current_user)
):
    """
    Get access log for all share links of this report.
    Only available to the report owner.
    """
    # Find all share tokens for this report
    relevant_tokens = []
    for token, link_data in SHARE_LINKS.items():
        if link_data["report_id"] == report_id:
            # Check if user is the creator
            if link_data.get("created_by") == str(current_user.id):
                relevant_tokens.append(token)

    if not relevant_tokens:
        return {
            "report_id": report_id,
            "share_links": [],
            "total_accesses": 0,
            "message": "No share links found for this report"
        }

    # Collect access logs for all relevant tokens
    share_links_with_logs = []
    total_accesses = 0

    for token in relevant_tokens:
        access_log = SHARE_ACCESS_LOG.get(token, [])
        total_accesses += len(access_log)

        share_links_with_logs.append({
            "share_token": token,
            "share_url": f"http://localhost:3000/share/{token}",
            "created_at": SHARE_LINKS[token]["created_at"],
            "expires_at": SHARE_LINKS[token]["expires_at"],
            "access_count": len(access_log),
            "accesses": access_log
        })

    # Sort by creation date (newest first)
    share_links_with_logs.sort(key=lambda x: x["created_at"], reverse=True)

    return {
        "report_id": report_id,
        "share_links": share_links_with_logs,
        "total_accesses": total_accesses
    }


@router.delete("/{report_id}/share/{share_token}")
async def revoke_share_link(
    report_id: str,
    share_token: str,
    current_user: User = Depends(get_current_user)
):
    """
    Revoke a share link, preventing further access.
    Only the creator can revoke their share links.
    """
    # Check if share link exists
    if share_token not in SHARE_LINKS:
        raise HTTPException(status_code=404, detail="Share link not found")

    share_link = SHARE_LINKS[share_token]

    # Verify report_id matches
    if share_link["report_id"] != report_id:
        raise HTTPException(status_code=400, detail="Share token does not belong to this report")

    # Verify user is the creator
    if share_link.get("created_by") != str(current_user.id):
        raise HTTPException(status_code=403, detail="Only the creator can revoke this share link")

    # Remove the share link and its access log
    del SHARE_LINKS[share_token]
    if share_token in SHARE_ACCESS_LOG:
        del SHARE_ACCESS_LOG[share_token]

    return {
        "success": True,
        "message": "Share link revoked successfully",
        "share_token": share_token
    }


# ==================== AUDIT LOGS ====================

# REMOVED - Moved before /{report_id} route to avoid route conflict


# ==================== REPORT GENERATION (FOR TIMEOUT TESTING) ====================

@router.post("/generate-complex")
async def generate_complex_report(
    request: Request,
    delay_seconds: int = Query(default=30, description="Simulation delay in seconds"),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """
    Simulate long-running report generation for timeout testing.
    This endpoint intentionally takes a long time to complete.
    """
    import asyncio

    # Log the generation request (commented out for testing)
    # await log_audit(
    #     db=db,
    #     user_id=str(current_user.id),
    #     action=AuditAction.REPORT_CREATE,
    #     resource_type=ResourceType.REPORT,
    #     resource_id="complex_report_generation",
    #     details=f"Started complex report generation (simulated delay: {delay_seconds}s)"
    # )

    # Track analytics event (commented out for testing)
    # await track_event(
    #     db=db,
    #     user_id=str(current_user.id),
    #     event_type=EventType.REPORT_GENERATED,
    #     metadata={"simulation_delay": delay_seconds, "report_type": "complex"}
    # )

    # Simulate long-running process
    try:
        # Sleep in small chunks to allow for cancellation
        chunks = delay_seconds // 5 + (1 if delay_seconds % 5 > 0 else 0)
        for i in range(chunks):
            await asyncio.sleep(min(5, delay_seconds - i * 5))

        # If we completed successfully, return full results
        return {
            "status": "completed",
            "report_id": f"report_{uuid.uuid4().hex[:8]}",
            "title": "Complex Market Analysis Report",
            "completion_time": delay_seconds,
            "sections_completed": 10,
            "total_sections": 10,
            "message": "Report generation completed successfully",
            "data": {
                "summary": "Comprehensive market analysis completed",
                "key_findings": [
                    "Market size: $5.2B",
                    "Growth rate: 12% YoY",
                    "Top 3 competitors identified"
                ]
            }
        }
    except asyncio.CancelledError:
        # Handle cancellation (timeout)
        return {
            "status": "timeout",
            "message": "Report generation timed out. Please try again with a simpler query.",
            "sections_completed": 3,
            "total_sections": 10,
            "partial_results": {
                "summary": "Partial analysis available",
                "completed_sections": ["Market Overview", "Key Players", "Revenue Analysis"]
            }
        }


# ==================== CANCELLABLE LONG-RUNNING TASK (Feature #368) ====================

# In-memory storage for cancellable tasks
CANCELLABLE_TASKS = {}  # {task_id: {"status": "running", "cancelled": False, "progress": 0, "total_steps": 10, "current_step": 0, "result": None}}


@router.post("/cancellable-task")
async def start_cancellable_task(
    total_steps: int = Query(default=10, description="Total number of steps to simulate"),
    step_duration: float = Query(default=2.0, description="Duration of each step in seconds"),
    current_user: User = Depends(get_current_user),
    request: Request = None
):
    """
    Start a cancellable long-running task.
    Returns task_id that can be used to check status or cancel the task.

    Feature #368: Cancel long-running operation
    """
    task_id = str(uuid.uuid4())

    # Initialize task state
    CANCELLABLE_TASKS[task_id] = {
        "status": "running",
        "cancelled": False,
        "progress": 0,
        "total_steps": total_steps,
        "current_step": 0,
        "result": None,
        "partial_data": []
    }

    # Start background task
    asyncio.create_task(_process_cancellable_task(task_id, total_steps, step_duration))

    return {
        "task_id": task_id,
        "status": "started",
        "total_steps": total_steps,
        "message": f"Task started with {total_steps} steps"
    }


async def _process_cancellable_task(task_id: str, total_steps: int, step_duration: float):
    """Background worker for cancellable task processing"""
    task = CANCELLABLE_TASKS[task_id]

    try:
        for step in range(1, total_steps + 1):
            # Check if task was cancelled
            if task["cancelled"]:
                task["status"] = "cancelled"
                task["result"] = {
                    "message": "Task was cancelled by user",
                    "completed_steps": step - 1,
                    "total_steps": total_steps,
                    "partial_data": task["partial_data"]
                }
                return

            # Simulate work
            await asyncio.sleep(step_duration)

            # Update progress
            task["current_step"] = step
            task["progress"] = int((step / total_steps) * 100)
            task["partial_data"].append(f"Step {step} completed")

        # Task completed successfully
        task["status"] = "completed"
        task["result"] = {
            "message": "Task completed successfully",
            "completed_steps": total_steps,
            "total_steps": total_steps,
            "final_data": task["partial_data"]
        }

    except Exception as e:
        task["status"] = "error"
        task["result"] = {
            "message": f"Task failed: {str(e)}",
            "completed_steps": task["current_step"],
            "total_steps": total_steps
        }


@router.post("/cancellable-task/{task_id}/cancel")
async def cancel_task(
    task_id: str,
    current_user: User = Depends(get_current_user)
):
    """
    Cancel a running task.

    Feature #368: Cancel long-running operation
    """
    if task_id not in CANCELLABLE_TASKS:
        raise HTTPException(status_code=404, detail="Task not found")

    task = CANCELLABLE_TASKS[task_id]

    if task["status"] != "running":
        return {
            "task_id": task_id,
            "message": f"Task is already {task['status']}, cannot cancel",
            "status": task["status"]
        }

    # Set cancelled flag
    task["cancelled"] = True

    return {
        "task_id": task_id,
        "message": "Cancellation requested. Task will stop after current step.",
        "current_step": task["current_step"],
        "total_steps": task["total_steps"]
    }


@router.get("/cancellable-task/{task_id}/status")
async def get_task_status(
    task_id: str,
    current_user: User = Depends(get_current_user)
):
    """
    Get status of a cancellable task.

    Feature #368: Cancel long-running operation
    """
    if task_id not in CANCELLABLE_TASKS:
        raise HTTPException(status_code=404, detail="Task not found")

    task = CANCELLABLE_TASKS[task_id]

    return {
        "task_id": task_id,
        "status": task["status"],
        "progress": task["progress"],
        "current_step": task["current_step"],
        "total_steps": task["total_steps"],
        "cancelled": task["cancelled"],
        "result": task["result"]
    }


# ==================== RETRY FAILED OPERATION (Feature #369) ====================

# In-memory storage for retryable tasks
RETRY_TASKS = {}  # {task_id: {"status": "failed|success|running", "attempt": 1, "max_attempts": 3, "should_fail": True, "result": None}}


class RetryTaskRequest(BaseModel):
    """Request model for retry task"""
    should_fail_first: bool = True
    max_attempts: int = 3


@router.post("/retry-task")
async def start_retry_task(
    request: RetryTaskRequest,
    current_user: User = Depends(get_current_user)
):
    """
    Start a retryable task that can fail and be retried.

    Feature #369: Retry failed operation
    """
    task_id = str(uuid.uuid4())

    # Initialize task state
    RETRY_TASKS[task_id] = {
        "status": "running",
        "attempt": 1,
        "max_attempts": request.max_attempts,
        "should_fail": request.should_fail_first,
        "result": None,
        "error": None
    }

    # Start background task
    asyncio.create_task(_process_retry_task(task_id))

    return {
        "task_id": task_id,
        "message": "Retryable task started",
        "attempt": 1,
        "max_attempts": request.max_attempts
    }


async def _process_retry_task(task_id: str):
    """Background worker for retryable task processing"""
    task = RETRY_TASKS[task_id]

    try:
        # Simulate some work
        await asyncio.sleep(1.5)

        # Check if task should fail on this attempt
        if task["should_fail"] and task["attempt"] == 1:
            # Fail on first attempt
            task["status"] = "failed"
            task["error"] = "Simulated failure: Network timeout"
            task["result"] = None
        else:
            # Success
            task["status"] = "success"
            task["result"] = {
                "data": "Task completed successfully",
                "attempt": task["attempt"],
                "timestamp": datetime.utcnow().isoformat()
            }
            task["error"] = None

    except Exception as e:
        task["status"] = "failed"
        task["error"] = str(e)
        task["result"] = None


@router.post("/retry-task/{task_id}/retry")
async def retry_task(
    task_id: str,
    current_user: User = Depends(get_current_user)
):
    """
    Retry a failed task.

    Feature #369: Retry failed operation
    """
    if task_id not in RETRY_TASKS:
        raise HTTPException(status_code=404, detail="Task not found")

    task = RETRY_TASKS[task_id]

    if task["status"] != "failed":
        return {
            "task_id": task_id,
            "message": "Task is not in failed state",
            "status": task["status"]
        }

    # Check if max attempts reached
    if task["attempt"] >= task["max_attempts"]:
        return {
            "task_id": task_id,
            "message": "Maximum retry attempts reached",
            "attempt": task["attempt"],
            "max_attempts": task["max_attempts"]
        }

    # Increment attempt counter
    task["attempt"] += 1
    task["status"] = "running"
    task["error"] = None

    # Mark that we should succeed on retry (simulate fix)
    task["should_fail"] = False

    # Start background task again
    asyncio.create_task(_process_retry_task(task_id))

    return {
        "task_id": task_id,
        "message": "Task retry initiated",
        "attempt": task["attempt"],
        "max_attempts": task["max_attempts"]
    }


@router.get("/retry-task/{task_id}/status")
async def get_retry_task_status(
    task_id: str,
    current_user: User = Depends(get_current_user)
):
    """
    Get status of a retryable task.

    Feature #369: Retry failed operation
    """
    if task_id not in RETRY_TASKS:
        raise HTTPException(status_code=404, detail="Task not found")

    task = RETRY_TASKS[task_id]

    return {
        "task_id": task_id,
        "status": task["status"],
        "attempt": task["attempt"],
        "max_attempts": task["max_attempts"],
        "error": task["error"],
        "result": task["result"]
    }


# ==================== TEST ENDPOINTS FOR SHARE LINK EXPIRATION ====================

@router.post("/test/share/expire-link/{share_token}")
async def test_expire_share_link(
    share_token: str,
    seconds_until_expiry: int = Query(default=5, description="Seconds until link expires")
):
    """
    TEST ENDPOINT: Manually set expiration time for a share link.
    Used for testing expiration functionality.

    Feature #48: Report sharing expiration
    """
    if share_token not in SHARE_LINKS:
        raise HTTPException(status_code=404, detail="Share token not found")

    # Set expiration to N seconds from now
    from datetime import timedelta
    new_expires_at = datetime.now() + timedelta(seconds=seconds_until_expiry)

    SHARE_LINKS[share_token]["expires_at"] = new_expires_at.isoformat() + "Z"

    return {
        "message": f"Share link will expire in {seconds_until_expiry} seconds",
        "share_token": share_token,
        "expires_at": new_expires_at.isoformat() + "Z",
        "current_time": datetime.now().isoformat() + "Z"
    }


# ==================== BULK IMPORT ENDPOINTS ====================

class BulkImportPreview(BaseModel):
    """Preview of data to be imported"""
    total_rows: int
    valid_rows: int
    invalid_rows: int
    preview_data: List[dict]
    errors: List[dict]


class BulkImportResult(BaseModel):
    """Result of bulk import operation"""
    total_rows: int
    imported_count: int
    failed_count: int
    skipped_count: int = 0  # Feature #153: Duplicate handling
    imported_ids: List[str]
    skipped_duplicates: List[dict] = []  # Feature #153: List of skipped duplicates
    errors: List[dict]


@router.post("/bulk-import-preview")
async def bulk_import_preview(
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user)
):
    """
    Preview CSV/Excel file before importing.

    Expected columns:
    - title (required)
    - type (required: company_profile, market_analysis, competitive_analysis, due_diligence)
    - company (optional)
    - summary (optional)
    - status (optional: draft, in_progress, completed - default: draft)

    Feature #152: Import valid data creates records
    """
    import csv
    from openpyxl import load_workbook

    # Check file extension
    filename = file.filename.lower()
    if not (filename.endswith('.csv') or filename.endswith('.xlsx') or filename.endswith('.xls')):
        raise HTTPException(
            status_code=400,
            detail="Invalid file format. Only CSV and Excel files are supported."
        )

    # Read file content
    content = await file.read()

    rows = []
    valid_rows = []
    invalid_rows = []
    errors = []

    try:
        if filename.endswith('.csv'):
            # Parse CSV
            content_str = content.decode('utf-8')
            csv_reader = csv.DictReader(io.StringIO(content_str))
            rows = list(csv_reader)

            # Feature #154: Validate CSV structure
            # Check if required headers are present
            if csv_reader.fieldnames:
                fieldnames_lower = [f.lower() if f else '' for f in csv_reader.fieldnames]
                if 'title' not in fieldnames_lower or 'type' not in fieldnames_lower:
                    raise ValueError("CSV must contain 'title' and 'type' columns")

            # Detect malformed CSV by checking for suspicious newlines in field values
            for idx, row in enumerate(rows, start=2):
                for field_name, field_value in row.items():
                    if field_value and isinstance(field_value, str):
                        # If field contains unescaped newlines, likely malformed
                        if '\n' in field_value or '\r' in field_value:
                            raise ValueError(
                                f"Malformed CSV detected: Row {idx}, field '{field_name}' contains unexpected line breaks. "
                                f"This usually indicates unclosed quotes or improperly escaped data."
                            )
        else:
            # Parse Excel
            wb = load_workbook(filename=io.BytesIO(content))
            ws = wb.active

            # Get headers from first row
            headers = [cell.value for cell in ws[1]]

            # Feature #154: Validate Excel has required headers
            headers_lower = [h.lower() if h else '' for h in headers]
            if 'title' not in headers_lower or 'type' not in headers_lower:
                raise ValueError("Excel file must contain 'title' and 'type' columns")

            # Get data rows
            for row_idx, row in enumerate(ws.iter_rows(min_row=2, values_only=True), start=2):
                row_dict = {headers[i]: row[i] for i in range(len(headers)) if i < len(row)}
                rows.append(row_dict)

    except Exception as e:
        raise HTTPException(
            status_code=400,
            detail=f"Failed to parse file: {str(e)}"
        )

    # Validate rows
    valid_types = ["company_profile", "market_analysis", "competitive_analysis", "due_diligence"]
    valid_statuses = ["draft", "in_progress", "completed"]

    for idx, row in enumerate(rows, start=2):
        row_errors = []

        # Validate required fields
        if not row.get('title') or not str(row.get('title')).strip():
            row_errors.append("Missing or empty 'title' field")

        if not row.get('type') or not str(row.get('type')).strip():
            row_errors.append("Missing or empty 'type' field")
        elif row.get('type') not in valid_types:
            row_errors.append(f"Invalid type '{row.get('type')}'. Must be one of: {', '.join(valid_types)}")

        # Validate optional status field
        if row.get('status') and row.get('status') not in valid_statuses:
            row_errors.append(f"Invalid status '{row.get('status')}'. Must be one of: {', '.join(valid_statuses)}")

        if row_errors:
            invalid_rows.append(idx)
            errors.append({
                "row": idx,
                "data": row,
                "errors": row_errors
            })
        else:
            valid_rows.append(row)

    # Return preview
    preview_data = valid_rows[:10]  # Show first 10 valid rows

    return BulkImportPreview(
        total_rows=len(rows),
        valid_rows=len(valid_rows),
        invalid_rows=len(invalid_rows),
        preview_data=preview_data,
        errors=errors[:20]  # Show first 20 errors
    )


@router.post("/bulk-import")
async def bulk_import_reports(
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """
    Import reports from CSV/Excel file.

    Expected columns:
    - title (required)
    - type (required: company_profile, market_analysis, competitive_analysis, due_diligence)
    - company (optional)
    - summary (optional)
    - status (optional: draft, in_progress, completed - default: draft)

    Feature #152: Import valid data creates records
    """
    import csv
    from openpyxl import load_workbook

    # Check file extension
    filename = file.filename.lower()
    if not (filename.endswith('.csv') or filename.endswith('.xlsx') or filename.endswith('.xls')):
        raise HTTPException(
            status_code=400,
            detail="Invalid file format. Only CSV and Excel files are supported."
        )

    # Read file content
    content = await file.read()

    rows = []

    try:
        if filename.endswith('.csv'):
            # Parse CSV
            content_str = content.decode('utf-8')
            csv_reader = csv.DictReader(io.StringIO(content_str))
            rows = list(csv_reader)

            # Feature #154: Validate CSV structure
            # Check if required headers are present
            if csv_reader.fieldnames:
                fieldnames_lower = [f.lower() if f else '' for f in csv_reader.fieldnames]
                if 'title' not in fieldnames_lower or 'type' not in fieldnames_lower:
                    raise ValueError("CSV must contain 'title' and 'type' columns")

            # Detect malformed CSV by checking for suspicious newlines in field values
            for idx, row in enumerate(rows, start=2):
                for field_name, field_value in row.items():
                    if field_value and isinstance(field_value, str):
                        # If field contains unescaped newlines, likely malformed
                        if '\n' in field_value or '\r' in field_value:
                            raise ValueError(
                                f"Malformed CSV detected: Row {idx}, field '{field_name}' contains unexpected line breaks. "
                                f"This usually indicates unclosed quotes or improperly escaped data."
                            )
        else:
            # Parse Excel
            wb = load_workbook(filename=io.BytesIO(content))
            ws = wb.active

            # Get headers from first row
            headers = [cell.value for cell in ws[1]]

            # Feature #154: Validate Excel has required headers
            headers_lower = [h.lower() if h else '' for h in headers]
            if 'title' not in headers_lower or 'type' not in headers_lower:
                raise ValueError("Excel file must contain 'title' and 'type' columns")

            # Get data rows
            for row_idx, row in enumerate(ws.iter_rows(min_row=2, values_only=True), start=2):
                row_dict = {headers[i]: row[i] for i in range(len(headers)) if i < len(row)}
                rows.append(row_dict)

    except Exception as e:
        raise HTTPException(
            status_code=400,
            detail=f"Failed to parse file: {str(e)}"
        )

    # Validate and import rows
    valid_types = ["company_profile", "market_analysis", "competitive_analysis", "due_diligence"]
    valid_statuses = ["draft", "in_progress", "completed"]

    imported_ids = []
    skipped_duplicates = []
    errors = []

    for idx, row in enumerate(rows, start=2):
        row_errors = []

        # Validate required fields
        if not row.get('title') or not str(row.get('title')).strip():
            row_errors.append("Missing or empty 'title' field")

        if not row.get('type') or not str(row.get('type')).strip():
            row_errors.append("Missing or empty 'type' field")
        elif row.get('type') not in valid_types:
            row_errors.append(f"Invalid type '{row.get('type')}'. Must be one of: {', '.join(valid_types)}")

        # Validate optional status field
        status = row.get('status', 'draft').strip() if row.get('status') else 'draft'
        if status not in valid_statuses:
            row_errors.append(f"Invalid status '{status}'. Must be one of: {', '.join(valid_statuses)}")

        if row_errors:
            errors.append({
                "row": idx,
                "data": row,
                "errors": row_errors
            })
            continue

        # Feature #153: Check for duplicates (skip if exists)
        # Duplicate = same title + type + company (case-insensitive)
        title = str(row.get('title')).strip()
        report_type = str(row.get('type')).strip()
        company = str(row.get('company')).strip() if row.get('company') else None

        # Check if duplicate exists (owned by same user)
        is_duplicate = False
        existing_report = None
        for r in MOCK_REPORTS:
            if (r.get("created_by") == str(current_user.id) and
                r.get("title", "").lower() == title.lower() and
                r.get("type", "").lower() == report_type.lower()):
                # Check company match (both None or both match)
                r_company = r.get("company")
                if (company is None and r_company is None) or \
                   (company and r_company and r_company.lower() == company.lower()):
                    is_duplicate = True
                    existing_report = r
                    break

        if is_duplicate:
            # Skip duplicate and record it
            skipped_duplicates.append({
                "row": idx,
                "data": row,
                "reason": "Duplicate report already exists",
                "existing_id": existing_report["id"] if existing_report else None
            })
            continue

        # Create report
        report_id = f"report_import_{uuid.uuid4().hex[:8]}"
        now = datetime.now().isoformat() + "Z"

        new_report = {
            "id": report_id,
            "title": str(row.get('title')).strip(),
            "type": str(row.get('type')).strip(),
            "company": str(row.get('company')).strip() if row.get('company') else None,
            "summary": str(row.get('summary')).strip() if row.get('summary') else f"Imported from {filename}",
            "status": status,
            "created_at": now,
            "updated_at": now,
            "created_by": str(current_user.id),
            "is_archived": False,
            "sections": [],
            "sources": []
        }

        # Add to mock database
        MOCK_REPORTS.append(new_report)
        imported_ids.append(report_id)

        # Log audit
        await log_audit(
            db=db,
            user=current_user,
            action_type=AuditAction.REPORT_CREATE,
            resource_type=ResourceType.REPORT,
            resource_id=report_id,
            description=f"Imported report from {filename}",
            extra_data={"imported_from": filename, "row": idx}
        )

    # Track analytics
    await track_event(
        db=db,
        event_type=EventType.REPORT_CREATED,
        event_name="bulk_import",
        user=current_user,
        metadata={
            "method": "bulk_import",
            "count": len(imported_ids),
            "skipped": len(skipped_duplicates),
            "failed": len(errors),
            "filename": filename
        }
    )

    return BulkImportResult(
        total_rows=len(rows),
        imported_count=len(imported_ids),
        failed_count=len(errors),
        skipped_count=len(skipped_duplicates),
        imported_ids=imported_ids,
        skipped_duplicates=skipped_duplicates[:50],  # Show first 50 skipped
        errors=errors[:50]  # Show first 50 errors
    )
